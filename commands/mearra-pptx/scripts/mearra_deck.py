"""
MearraDeck — high-level builder for brand-perfect Mearra presentations.

Philosophy:
  - The official Mearra template (assets/template.pptx) is the single source of
    truth for colors, fonts, logo, tagline, and layout grids.
  - This builder opens the template, strips the example slides, and provides
    a small, legible API to add on-brand slides by layout name.
  - It does NOT re-invent the grid — it picks layouts from the master by name
    and fills their placeholders.
  - For custom slides, it exposes `add_custom_slide()` which starts from a
    blank TITLE_ONLY layout (which still carries the master's brand chrome)
    and lets you paint on top with brand-safe shapes.

Usage:
    from scripts.mearra_deck import MearraDeck
    deck = MearraDeck()
    deck.add_cover(title="Wunder is now Mearra", subtitle="...")
    deck.add_section_header("Agenda")
    deck.add_title_body(title="...", body="...")
    deck.save("/path/to/out.pptx")

Every method returns the underlying python-pptx `Slide` object so you can
tweak it further (e.g., add a picture, a shape, a note) when needed.
"""
from __future__ import annotations

import os
import copy
from pathlib import Path
from typing import Iterable, Optional, List, Dict, Any

from pptx import Presentation
from pptx.slide import Slide
from pptx.util import Emu, Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.chart import XL_CHART_TYPE, XL_LABEL_POSITION, XL_LEGEND_POSITION
from pptx.chart.data import CategoryChartData
from pptx.oxml.ns import qn
from lxml import etree
import math

from .colors import (
    BLOSSOM, CORAL, INDIGO, TRENCH, LAGOON, CURRENT,
    MIST, WHITE, BLACK,
    FONT_HEADING, FONT_BODY,
    SLOGAN, URL, COMPANY, ALL_BRAND_COLORS,
)

HERE = Path(__file__).resolve().parent
SKILL_ROOT = HERE.parent
TEMPLATE = SKILL_ROOT / "assets" / "template.pptx"
ASSETS = SKILL_ROOT / "assets"


def _initials(name: str, max_chars: int = 2) -> str:
    """First letter of each word, up to `max_chars`. 'Jean Duckling' -> 'JD'."""
    words = [w for w in (name or "").split() if w]
    return "".join(w[0].upper() for w in words[:max_chars]) or "M"


def _split_bold(text: str):
    """
    Split `text` into (segment, is_bold) tokens based on `**…**` markers.
    Unclosed markers are treated as literal asterisks.
    """
    out = []
    i = 0
    bold = False
    buf = []
    while i < len(text):
        if text.startswith("**", i):
            if buf:
                out.append(("".join(buf), bold))
                buf = []
            bold = not bold
            i += 2
        else:
            buf.append(text[i])
            i += 1
    if buf:
        out.append(("".join(buf), bold))
    return out or [("", False)]


def _sample_brand_gradient(n: int):
    """
    Return `n` brand RGBColor values across an approved warm→cool gradient.
    Used for pyramid tiers / similar multi-band charts.
    """
    # Hard-coded ordered palette walk — all colours strictly in the brand bible.
    from .colors import BLOSSOM, CORAL, CURRENT, LAGOON, INDIGO
    seq_3 = [BLOSSOM, CORAL, INDIGO]
    seq_4 = [BLOSSOM, CORAL, LAGOON, INDIGO]
    seq_5 = [BLOSSOM, CORAL, CURRENT, LAGOON, INDIGO]
    return {3: seq_3, 4: seq_4, 5: seq_5}.get(n, seq_5[:n])


def _indigo_tints(n: int):
    """
    Return `n` RGBColor values interpolated from a pale indigo wash to full INDIGO.

    Mearra brand slides 28–31 use banded indigo tints to differentiate chart
    segments while staying on-brand. Pale end is a soft violet that sits
    comfortably on white; saturated end is INDIGO (#4400F0).
    """
    from pptx.dml.color import RGBColor
    # Pale end: soft violet derived from INDIGO at ~15% saturation on white.
    # Saturated end: INDIGO (#4400F0) = (68, 0, 240).
    r0, g0, b0 = 226, 220, 255   # very pale violet
    r1, g1, b1 = 68, 0, 240      # INDIGO
    if n <= 1:
        return [RGBColor(r1, g1, b1)]
    out = []
    for i in range(n):
        t = i / (n - 1)
        r = round(r0 + (r1 - r0) * t)
        g = round(g0 + (g1 - g0) * t)
        b = round(b0 + (b1 - b0) * t)
        out.append(RGBColor(r, g, b))
    return out


def _is_dark(color) -> bool:
    """
    Return True if `color` (RGBColor) is dark enough that white/MIST text reads
    better on top of it than dark TRENCH text.

    Uses a simple perceived-luminance heuristic (Rec. 709 coefficients).
    """
    try:
        r, g, b = color[0], color[1], color[2]
    except Exception:
        # RGBColor supports indexing; fall back to treating unknown as light.
        return False
    # Rec. 709 luma, normalised 0..1
    luma = (0.2126 * r + 0.7152 * g + 0.0722 * b) / 255.0
    return luma < 0.55


class MearraDeck:
    """Builder for brand-perfect Mearra presentations."""

    def __init__(self, template_path: Optional[str | Path] = None, strip_examples: bool = True):
        self.path = Path(template_path) if template_path else TEMPLATE
        if not self.path.exists():
            raise FileNotFoundError(
                f"Mearra template not found at {self.path}. "
                f"This skill must be installed with assets/template.pptx."
            )
        self.prs = Presentation(str(self.path))
        if strip_examples:
            self._strip_all_slides()
        # Build lookup: layout name -> SlideLayout
        self.layouts: Dict[str, Any] = {
            L.name: L for L in self.prs.slide_masters[0].slide_layouts
        }

    # ------------------------------------------------------------------ utilities

    def _strip_all_slides(self) -> None:
        """Remove every example slide from the template so we start fresh."""
        sldIdLst = self.prs.slides._sldIdLst
        r_ns = "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id"
        prs_part = self.prs.part
        for sldId in list(sldIdLst):
            rId = sldId.get(r_ns)
            # Drop the relationship from the presentation part pointing at the slide part.
            if rId is not None:
                try:
                    prs_part.drop_rel(rId)
                except KeyError:
                    pass
            sldIdLst.remove(sldId)

    def get_layout(self, name: str):
        """Return the SlideLayout with the given name. Raises if missing."""
        if name not in self.layouts:
            available = ", ".join(sorted(self.layouts.keys()))
            raise KeyError(f"No layout named {name!r}. Available: {available}")
        return self.layouts[name]

    def _add_slide(self, layout_name: str) -> Slide:
        return self.prs.slides.add_slide(self.get_layout(layout_name))

    @staticmethod
    def _fill_placeholder(slide: Slide, idx: int, text: str) -> None:
        """Set the text of a placeholder identified by its idx. Silently skips if missing."""
        for ph in slide.placeholders:
            if ph.placeholder_format.idx == idx:
                tf = ph.text_frame
                tf.text = text
                return
        # Not found — do nothing (some layouts omit certain placeholders)

    @staticmethod
    def _fill_placeholders(slide: Slide, mapping: Dict[int, str]) -> None:
        for idx, txt in mapping.items():
            if txt is None:
                continue
            MearraDeck._fill_placeholder(slide, idx, txt)

    @staticmethod
    def _suppress_placeholder_bullets(slide: Slide, idx: int) -> None:
        """
        Remove master-inherited bullets from every paragraph in the placeholder at `idx`.
        Useful for body placeholders on multi-column layouts where the master applies a
        default bullet character but we want clean prose.
        """
        for ph in slide.placeholders:
            if ph.placeholder_format.idx != idx:
                continue
            for p in ph.text_frame.paragraphs:
                pPr = p._p.get_or_add_pPr()
                for tag in ("a:buChar", "a:buAutoNum", "a:buNone"):
                    for el in pPr.findall(qn(tag)):
                        pPr.remove(el)
                etree.SubElement(pPr, qn("a:buNone"))
            return

    @staticmethod
    def _fill_heading_body(slide: Slide, idx: int, heading: str, body: str) -> None:
        """
        Fill a body placeholder with a bold heading paragraph followed by a regular body paragraph.
        Inherits the placeholder's font and colour from the layout/master — so the heading
        just gets bold, not recolored.

        The layout's body placeholders carry default bullet formatting inherited from the
        master; we explicitly suppress bullets on both paragraphs so the heading reads as a
        standalone title, not a sibling bullet to the body.
        """
        def _suppress_bullet(p) -> None:
            pPr = p._pPr if p._pPr is not None else p._p.get_or_add_pPr()
            # Remove any pre-existing bullet character/auto-num element and add <a:buNone/>.
            for tag in ("a:buChar", "a:buAutoNum", "a:buNone"):
                for el in pPr.findall(qn(tag)):
                    pPr.remove(el)
            etree.SubElement(pPr, qn("a:buNone"))

        for ph in slide.placeholders:
            if ph.placeholder_format.idx != idx:
                continue
            tf = ph.text_frame
            tf.clear()
            p1 = tf.paragraphs[0]
            _suppress_bullet(p1)
            r1 = p1.add_run()
            r1.text = heading
            r1.font.bold = True
            p2 = tf.add_paragraph()
            _suppress_bullet(p2)
            r2 = p2.add_run()
            r2.text = body
            return

    # ---- Decorative accent placement ----------------------------------------

    # Accent style presets — position + size (inches), slide is 10.00 × 5.62.
    # Each preset lists an ordered fallback of real Aquatic Form clusters suitable for
    # that composition. The first form in the list is the default for 'auto'.
    #
    # Rule of thumb: only reference the four genuine Aquatic Form clusters (tall-01..04).
    # cluster-full-02 is a soft gradient orb and reads as "just an ellipse" on its own —
    # it's kept available as a background primitive but is NOT used as a default accent.
    _ACCENT_STYLES: Dict[str, Dict[str, Any]] = {
        # Thin tall accent bleeding off the right edge — for quote / metric / hero slides
        # that benefit from a warm multi-colour cluster anchor.
        "tall-right": {
            "left": 7.6, "top": -0.2, "width": 2.8, "height": 5.9,
            "forms": ["cluster-tall-01", "cluster-tall-03", "cluster-tall-02", "cluster-tall-04"],
        },
        # Mirror of tall-right, bleeds off the left.
        "tall-left": {
            "left": -0.4, "top": -0.2, "width": 2.8, "height": 5.9,
            "forms": ["cluster-tall-02", "cluster-tall-04", "cluster-tall-03", "cluster-tall-01"],
        },
        # Bottom-right anchor using a real tall cluster (not the soft orb) — for when a
        # slide explicitly asks for a decorative corner anchor.
        "bottom-right": {
            "left": 7.2, "top": 1.7, "width": 3.2, "height": 4.2,
            "forms": ["cluster-tall-03", "cluster-tall-01"],
        },
    }

    def _place_auto_accent(
        self,
        slide: Slide,
        style: str = "tall-right",
        requested: Optional[str] = "auto",
    ):
        """
        Place an Aquatic Form as a decorative accent on `slide`.

        style: one of `_ACCENT_STYLES` keys — controls position + size.
        requested:
          * "auto" — use the first form in the style's form list (on-brand default).
          * "<name>" — use a specific aquatic form (e.g. "blob-dark-01"). Fuzzy matched.
          * None / "" / "none" — do nothing.

        The aquatic form is inserted BEHIND any placeholders so text stays readable —
        we do this by moving the shape to the start of the slide's shape tree.
        """
        if not requested or requested in ("none", "off"):
            return None
        preset = self._ACCENT_STYLES.get(style)
        if preset is None:
            raise ValueError(
                f"Unknown accent style {style!r}. Choose from: {list(self._ACCENT_STYLES)}"
            )

        # Resolve which form to use.
        if requested == "auto":
            form_name = preset["forms"][0]
        else:
            form_name = requested

        try:
            pic = self.add_aquatic_form(
                slide, name=form_name,
                left=preset["left"], top=preset["top"],
                width=preset["width"], height=preset["height"],
            )
        except FileNotFoundError:
            # Fall back to any form in the style's list.
            for fallback in preset["forms"]:
                try:
                    pic = self.add_aquatic_form(
                        slide, name=fallback,
                        left=preset["left"], top=preset["top"],
                        width=preset["width"], height=preset["height"],
                    )
                    break
                except FileNotFoundError:
                    continue
            else:
                return None

        # Push the picture to the bottom of the z-order so text placeholders sit on top.
        self._send_to_back(pic)
        return pic

    @staticmethod
    def _resize_placeholder(
        slide,
        idx: int,
        left: Optional[float] = None,
        top: Optional[float] = None,
        width: Optional[float] = None,
        height: Optional[float] = None,
    ) -> None:
        """
        Explicitly set the placeholder's geometry (inches). When python-pptx writes
        any single dimension on an inherited placeholder it drops the others to 0,
        so we always read current values first and write all four together.
        """
        for ph in slide.placeholders:
            if ph.placeholder_format.idx != idx:
                continue
            cur_l = ph.left
            cur_t = ph.top
            cur_w = ph.width
            cur_h = ph.height
            ph.left = Inches(left) if left is not None else cur_l
            ph.top = Inches(top) if top is not None else cur_t
            ph.width = Inches(width) if width is not None else cur_w
            ph.height = Inches(height) if height is not None else cur_h
            return

    @classmethod
    def _shift_placeholder(cls, slide, idx: int, dy_inches: float) -> None:
        """Move placeholder `idx` down by `dy_inches` (negative = up)."""
        for ph in slide.placeholders:
            if ph.placeholder_format.idx == idx:
                new_top = (ph.top or 0) / 914400 + dy_inches
                cls._resize_placeholder(slide, idx, top=new_top)
                return

    @classmethod
    def _shrink_placeholder(cls, slide, idx: int, dy_inches: float) -> None:
        """Shrink placeholder `idx` in height by `dy_inches`."""
        for ph in slide.placeholders:
            if ph.placeholder_format.idx == idx:
                new_h = (ph.height or 0) / 914400 - dy_inches
                cls._resize_placeholder(slide, idx, height=new_h)
                return

    @staticmethod
    def _send_to_back(shape) -> None:
        """Move `shape` to the first child position in its parent spTree (back of z-order)."""
        sp = shape._element
        parent = sp.getparent()
        if parent is None:
            return
        parent.remove(sp)
        # spTree's first two children are nvGrpSpPr and grpSpPr — insert after those.
        # Safest: insert at index 2 so the accent sits at the back of actual shapes.
        insert_at = 2 if len(parent) >= 2 else 0
        parent.insert(insert_at, sp)

    # ------------------------------------------------------------------ slide types

    def add_cover(
        self,
        title: str,
        subtitle: Optional[str] = None,
        presenter_line: Optional[str] = None,
        variant: str = "indigo",
    ) -> Slide:
        """
        Cover / title slide.

        variant:
          * "indigo" — full-bleed dark indigo background with Aquatic Forms cluster on the right
            (the most common opening slide — uses layout 'TITLE')
          * "light" — Mist background with Aquatic Forms cluster on the right
            (uses layout 'SECTION_HEADER' which serves as the light cover in this template)

        The master slide handles the tagline "Business. People. Tech. Together."
        and logo placement automatically.
        """
        if variant == "indigo":
            slide = self._add_slide("TITLE")
        elif variant == "light":
            slide = self._add_slide("SECTION_HEADER")
        else:
            raise ValueError("variant must be 'indigo' or 'light'")
        # Placeholder idx=0 is the title, idx=1 is a subtitle line,
        # idx=2 is a secondary subtitle (used for presenter/date)
        self._fill_placeholders(slide, {0: title, 1: subtitle or "", 2: presenter_line or ""})
        return slide

    def add_section_header(self, title: str, eyebrow: Optional[str] = None) -> Slide:
        """
        Full-bleed Indigo section divider (e.g., 'Agenda', 'Mearra Method', 'Concept case study').

        Uses a blue-background custom layout — the title is big italic white.
        If `eyebrow` is given, it appears as the small supporting line.
        """
        slide = self._add_slide("CUSTOM_8_1_2_1_1")
        # This layout has title at idx=0 and 5 subtitle slots (idx 1..5).
        self._fill_placeholders(slide, {0: title, 1: eyebrow or ""})
        return slide

    def add_title_body(
        self,
        title: str,
        body: str,
        eyebrow: Optional[str] = None,
        accent: Optional[str] = None,
    ) -> Slide:
        """
        Mist background, title + single body column. `eyebrow` is the tiny label top-left.

        In the Mearra brand, concept slides are typically clean text on Mist with no
        decoration — the master chrome (logo, tagline) carries the branding. Accents are
        opt-in for moments that need extra visual punch, not the default.

        `accent` when set places a subtle Aquatic Form decoration on the right-bleed:
        - None (default) → clean text slide, matches brand originals
        - "auto"         → right-edge tall cluster (warm, multi-color)
        - "<form-name>"  → specific aquatic form from assets/aquatic-forms/
        """
        slide = self._add_slide("TITLE_AND_BODY")
        if accent:
            # Reserve right-side real estate for the accent so text doesn't cross it.
            self._resize_placeholder(slide, 0, width=6.8)
            self._resize_placeholder(slide, 1, width=6.8)
        self._fill_placeholders(slide, {0: title, 1: body, 2: eyebrow or ""})
        self._suppress_placeholder_bullets(slide, 1)
        if accent:
            self._place_auto_accent(slide, style="tall-right", requested=accent)
        return slide

    def add_title_subtitle_body(
        self,
        title: str,
        subtitle: str,
        body: str,
        eyebrow: Optional[str] = None,
        accent: Optional[str] = None,
    ) -> Slide:
        """
        Mist background, title + subtitle + single body column.

        Like `add_title_body`, brand originals are clean — no accent by default. Pass
        `accent="auto"` (or a specific form name) only when a slide genuinely benefits
        from the extra visual weight.
        """
        slide = self._add_slide("TITLE_AND_BODY_2")
        if accent:
            self._resize_placeholder(slide, 0, width=6.8)
            self._resize_placeholder(slide, 1, width=6.8)
            self._resize_placeholder(slide, 3, width=6.8)
        self._fill_placeholders(slide, {0: title, 1: body, 2: eyebrow or "", 3: subtitle})
        self._suppress_placeholder_bullets(slide, 1)
        if accent:
            self._place_auto_accent(slide, style="tall-right", requested=accent)
        return slide

    def add_title_only(
        self,
        title: str,
        eyebrow: Optional[str] = None,
        accent: Optional[str] = None,
    ) -> Slide:
        """
        Statement slide — one strong headline on Mist, nothing else.

        The brand uses these as declarative moments. Keep them clean — no gradient ellipses,
        no corner graphics. Pass `accent="auto"` only when the statement genuinely needs
        a visual punch (rare — typically a section-opener does better).
        """
        slide = self._add_slide("TITLE_ONLY")
        self._fill_placeholders(slide, {0: title, 2: eyebrow or ""})
        if accent:
            self._place_auto_accent(slide, style="tall-right", requested=accent)
        return slide

    def add_title_subtitle_only(
        self,
        title: str,
        subtitle: str,
        eyebrow: Optional[str] = None,
        accent: Optional[str] = None,
    ) -> Slide:
        """
        Chapter opener — title + one-line subtitle on Mist, nothing else.

        Matches the brand originals (slide-09, slide-13): clean text anchored top-left,
        plenty of whitespace. If the chapter wants visual weight, prefer `add_section_header`
        (full-bleed indigo) or `add_hero` (full-bleed photo) over decorating this one.
        """
        slide = self._add_slide("TITLE_ONLY_2")
        self._fill_placeholders(slide, {0: title, 2: eyebrow or "", 3: subtitle})
        if accent:
            self._place_auto_accent(slide, style="tall-right", requested=accent)
        return slide

    def add_two_columns(
        self,
        title: str,
        left_title: str,
        left_body: str,
        right_title: str,
        right_body: str,
        eyebrow: Optional[str] = None,
        left_icon: Optional[str] = None,
        right_icon: Optional[str] = None,
    ) -> Slide:
        """
        Title + two side-by-side columns, each with its own bold heading and body paragraph.
        Uses the TITLE_AND_TWO_COLUMNS layout which has exactly 2 body columns; the column
        heading is rendered as a bold first paragraph in the body placeholder.

        Optional `left_icon` / `right_icon` names an icon from assets/icons/ to place
        above each column's heading (small, 0.65in square).
        """
        slide = self._add_slide("TITLE_AND_TWO_COLUMNS")
        # idx 0 = title, 3 = eyebrow, 1 = L body, 2 = R body
        self._fill_placeholders(slide, {0: title, 3: eyebrow or ""})
        # Body placeholders live at top=1.70. If we have icons, push them down
        # to 2.55 (0.85" lower) and shrink height so they still end at the same bottom.
        if left_icon or right_icon:
            for body_idx in (1, 2):
                self._shift_placeholder(slide, body_idx, 0.85)
                self._shrink_placeholder(slide, body_idx, 0.85)
        self._fill_heading_body(slide, 1, left_title,  left_body)
        self._fill_heading_body(slide, 2, right_title, right_body)
        # Icons sit in the freed row — aligned with their body column.
        if left_icon:
            self.add_icon(slide, left_icon, left=0.34, top=1.80, size=0.65)
        if right_icon:
            self.add_icon(slide, right_icon, left=5.38, top=1.80, size=0.65)
        return slide

    def add_three_columns(
        self,
        title: str,
        columns: Iterable[Dict[str, str]],
        eyebrow: Optional[str] = None,
    ) -> Slide:
        """
        Title + three columns. `columns` is an iterable of dicts:
            [{'title': '...', 'body': '...', 'icon': '<optional icon name>'}, ...]
        (exactly 3 columns). `icon` is optional and resolves against assets/icons/.
        """
        cols = list(columns)
        if len(cols) != 3:
            raise ValueError("add_three_columns requires exactly 3 columns")
        slide = self._add_slide("TITLE_AND_TWO_COLUMNS_1_1_1")
        # Layout positions: heading @ top=1.55, body @ top=2.00. If any icons,
        # push both placeholders down by 0.80" and shrink body height accordingly.
        has_icons = any(c.get("icon") for c in cols)
        if has_icons:
            for idx in (3, 6, 7):  # column headings
                self._shift_placeholder(slide, idx, 0.80)
            for idx in (1, 4, 5):  # column bodies
                self._shift_placeholder(slide, idx, 0.80)
                self._shrink_placeholder(slide, idx, 0.80)
        # idx 0 = title, 2 = eyebrow, 3/6/7 = column titles, 1/4/5 = column bodies
        self._fill_placeholders(slide, {
            0: title, 2: eyebrow or "",
            3: cols[0]["title"], 1: cols[0]["body"],
            6: cols[1]["title"], 4: cols[1]["body"],
            7: cols[2]["title"], 5: cols[2]["body"],
        })
        # Master applies bullets to body placeholders — suppress them for clean prose.
        for body_idx in (1, 4, 5):
            self._suppress_placeholder_bullets(slide, body_idx)
        # Icons go in the freed row — aligned with the left edge of each column.
        icon_lefts = [0.34, 3.56, 6.78]
        for col, lx in zip(cols, icon_lefts):
            if col.get("icon"):
                self.add_icon(slide, col["icon"], left=lx, top=1.55, size=0.65)
        return slide

    def add_four_columns(
        self,
        title: str,
        columns: Iterable[Dict[str, str]],
        eyebrow: Optional[str] = None,
    ) -> Slide:
        """Title + four columns. Each column dict supports `title`, `body`, optional `icon`."""
        cols = list(columns)
        if len(cols) != 4:
            raise ValueError("add_four_columns requires exactly 4 columns")
        slide = self._add_slide("TITLE_AND_TWO_COLUMNS_1_1_1_2")
        has_icons = any(c.get("icon") for c in cols)
        if has_icons:
            for idx in (3, 6, 7, 9):  # column headings @ y=1.55
                self._shift_placeholder(slide, idx, 0.75)
            for idx in (1, 4, 5, 8):  # column bodies @ y=2.15
                self._shift_placeholder(slide, idx, 0.75)
                self._shrink_placeholder(slide, idx, 0.75)
        self._fill_placeholders(slide, {
            0: title, 2: eyebrow or "",
            3: cols[0]["title"], 1: cols[0]["body"],
            6: cols[1]["title"], 4: cols[1]["body"],
            7: cols[2]["title"], 5: cols[2]["body"],
            9: cols[3]["title"], 8: cols[3]["body"],
        })
        # Master applies bullets to body placeholders — suppress them for clean prose.
        for body_idx in (1, 4, 5, 8):
            self._suppress_placeholder_bullets(slide, body_idx)
        # Four column layout has narrower columns at 0.34, 2.77, 5.20, 7.63.
        icon_lefts = [0.34, 2.77, 5.20, 7.63]
        for col, lx in zip(cols, icon_lefts):
            if col.get("icon"):
                self.add_icon(slide, col["icon"], left=lx, top=1.55, size=0.55)
        return slide

    # ---- New high-impact slide types ----------------------------------------

    def add_big_number(
        self,
        number: str,
        label: str,
        context: Optional[str] = None,
        eyebrow: Optional[str] = None,
        accent: Optional[str] = "auto",
    ) -> Slide:
        """
        Statement-style metric slide: giant number + one-line label + optional context.

        E.g. `add_big_number("15+", "years dismantling complexity", context="…for Finland's largest media houses")`

        Built on the TITLE_ONLY layout (blank canvas with brand chrome) so we can place
        the number ourselves in Sofia Sans Extra Condensed Black Italic — the signature
        Mearra statement style.
        """
        slide = self._add_slide("TITLE_ONLY")
        # Clear the title placeholder (it would otherwise dominate).
        self._fill_placeholder(slide, 0, "")
        if eyebrow:
            self._fill_placeholder(slide, 2, eyebrow)

        # Giant number, left-anchored. Sized to leave room for label + context beneath.
        self.add_heading_textbox(
            slide, number,
            left=0.6, top=1.0, width=6.0, height=2.6,
            color=INDIGO, size_pt=180, italic=True,
        )
        # Label — bold body, Trench, right under the number.
        self.add_body_textbox(
            slide, label,
            left=0.7, top=3.75, width=6.0, height=0.5,
            color=TRENCH, size_pt=22, bold=True,
        )
        if context:
            self.add_body_textbox(
                slide, context,
                left=0.7, top=4.3, width=6.0, height=1.0,
                color=TRENCH, size_pt=14,
            )
        if accent:
            self._place_auto_accent(slide, style="tall-right", requested=accent)
        return slide

    def add_metric_strip(
        self,
        title: str,
        metrics: Iterable[Dict[str, str]],
        eyebrow: Optional[str] = None,
    ) -> Slide:
        """
        Title + horizontal strip of 3–4 metric blocks.
        metrics = [{"number": "5", "label": "new clients"}, ...]

        Numbers are rendered in Sofia Sans Extra Condensed Black Italic Indigo,
        labels in Sofia Sans Medium Trench. Great for "Q1 wins" / "by the numbers" slides.
        """
        metrics = list(metrics)
        if not 2 <= len(metrics) <= 4:
            raise ValueError("add_metric_strip needs between 2 and 4 metrics")
        slide = self._add_slide("TITLE_ONLY")
        self._fill_placeholder(slide, 0, title)
        if eyebrow:
            self._fill_placeholder(slide, 2, eyebrow)

        # Strip sits in the bottom 2/3 of the slide. Compute widths dynamically.
        n = len(metrics)
        margin = 0.55
        gap = 0.35
        usable_w = 10.0 - 2 * margin - (n - 1) * gap
        block_w = usable_w / n
        top_num = 2.6
        for i, m in enumerate(metrics):
            left = margin + i * (block_w + gap)
            # Number
            self.add_heading_textbox(
                slide, m.get("number", ""),
                left=left, top=top_num, width=block_w, height=1.4,
                color=INDIGO, size_pt=88, italic=True,
            )
            # Label
            self.add_body_textbox(
                slide, m.get("label", ""),
                left=left, top=top_num + 1.4, width=block_w, height=0.8,
                color=TRENCH, size_pt=14, bold=True,
            )
        return slide

    def add_hero(
        self,
        title: str,
        subtitle: Optional[str] = None,
        photo: Optional[str] = None,
        eyebrow: Optional[str] = None,
        accent: Optional[str] = None,
    ) -> Slide:
        """
        Full-bleed hero slide — big title anchored bottom-left over a background photo
        (or a full-bleed aquatic cluster). Use as a chapter opener or emotional punctuation.

        `photo`  : name from assets/backgrounds/ (e.g. "atmospheric-orange") — full-bleed.
        `accent` : alternative to photo — a full-bleed aquatic cluster. Photo wins if both.

        Title sits in a 48pt italic block anchored just above the bottom margin; subtitle
        sits directly beneath it with a 0.15" gap so the two never collide.
        """
        slide = self._add_slide("TITLE_ONLY")
        self._fill_placeholder(slide, 0, "")  # clear master title so we draw our own
        if eyebrow:
            self._fill_placeholder(slide, 2, eyebrow)

        dark_bg = False
        if photo:
            try:
                img_path = self._find_asset("backgrounds", photo, prefix="photo-")
                pic = self.add_image(slide, img_path, left=0, top=0, width=10.0, height=5.625)
                self._send_to_back(pic)
                dark_bg = True
            except FileNotFoundError:
                pass
        elif accent:
            # Full-bleed aquatic cluster — place directly (no preset needed).
            form_name = accent if accent != "auto" else "cluster-tall-01"
            try:
                pic = self.add_aquatic_form(
                    slide, name=form_name,
                    left=0.0, top=0.0, width=10.0, height=5.625,
                )
                self._send_to_back(pic)
                dark_bg = True
            except FileNotFoundError:
                pass

        # Title: 48pt italic, 3.0" wide x 1.2" tall block anchored at y=3.05 (bottom ≈ 4.25)
        # Subtitle (if any): 16pt, starts at y=4.40 with 0.15" gap below title → no collision.
        title_color = WHITE if dark_bg else TRENCH
        self.add_heading_textbox(
            slide, title,
            left=0.55, top=3.05, width=8.8, height=1.2,
            color=title_color, size_pt=48, italic=True,
        )
        if subtitle:
            self.add_body_textbox(
                slide, subtitle,
                left=0.6, top=4.40, width=8.8, height=0.55,
                color=title_color, size_pt=16, bold=True,
            )
        return slide

    def add_quote(
        self,
        quote: str,
        attribution: Optional[str] = None,
        eyebrow: Optional[str] = None,
        accent: Optional[str] = "auto",
    ) -> Slide:
        """
        Pull-quote slide — oversized italic quote, small attribution below.
        Use for testimonials or founding-principle lines.
        """
        slide = self._add_slide("TITLE_ONLY")
        self._fill_placeholder(slide, 0, "")
        if eyebrow:
            self._fill_placeholder(slide, 2, eyebrow)

        self.add_heading_textbox(
            slide, f"\u201c{quote}\u201d",
            left=0.7, top=1.1, width=7.0, height=3.2,
            color=TRENCH, size_pt=40, italic=True,
        )
        if attribution:
            self.add_body_textbox(
                slide, f"\u2014 {attribution}",
                left=0.75, top=4.4, width=7.0, height=0.6,
                color=TRENCH, size_pt=14, bold=True,
            )
        if accent:
            self._place_auto_accent(slide, style="tall-right", requested=accent)
        return slide

    # ---- Split image + text ------------------------------------------------

    def add_split_image_text(
        self,
        title: str,
        body: str,
        image: Optional[str] = None,
        eyebrow: Optional[str] = None,
        image_side: str = "left",
        accent_color: RGBColor = CORAL,
    ) -> Slide:
        """
        Split composition — half of the slide is a photo/aquatic panel, the other half
        holds the title + body text on Mist. Modelled on brand slide-38.

        `image`      : background name from assets/backgrounds/ (e.g. "atmospheric-orange").
                       If missing, a solid brand block (Coral or accent_color) is used
                       with a corner aquatic accent so the slide still reads on-brand.
        `image_side` : "left" (default) or "right".
        `accent_color`: colour for the solid fallback block when no image is provided.

        Text sits on the non-image half with comfortable padding. If the body contains
        a double-newline, each paragraph is emitted separately; single newlines are
        preserved as line breaks. Bold phrasing uses the same `**word**` convention as
        `add_rich_body` so you can stress words the Mearra way.
        """
        if image_side not in ("left", "right"):
            raise ValueError("image_side must be 'left' or 'right'")
        slide = self._add_slide("TITLE_ONLY")
        self._fill_placeholder(slide, 0, "")  # clear master title so we draw our own

        # Geometry: image panel is 4.2" wide; text gets the rest with 0.55" outer padding.
        panel_w = 4.2
        if image_side == "left":
            panel_left = 0.0
            text_left = panel_w + 0.55
            text_width = 10.0 - text_left - 0.55
        else:
            panel_left = 10.0 - panel_w
            text_left = 0.55
            text_width = panel_left - text_left - 0.35

        # Image panel — photo full-bleed, or solid fallback + aquatic ghost.
        # Look across portraits, mockups, and backgrounds so the caller can pass
        # any image-family name without worrying about which folder it lives in.
        img_path = None
        if image:
            for folder, prefix in (
                ("portraits", "person-"),
                ("mockups", None),
                ("backgrounds", "photo-"),
                ("graphics", None),
            ):
                try:
                    img_path = self._find_asset(folder, image, prefix=prefix)
                    break
                except FileNotFoundError:
                    continue
        if img_path:
            pic = self.add_image(
                slide, img_path, left=panel_left, top=0.0,
                width=panel_w, height=5.625,
                fit="cover",  # preserve aspect; centre-crop overflow
            )
            self._send_to_back(pic)
        else:
            self._paint_solid_panel(slide, panel_left, 0.0, panel_w, 5.625, accent_color)

        # Text column
        if eyebrow:
            self.add_body_textbox(
                slide, eyebrow,
                left=text_left, top=0.5, width=text_width, height=0.4,
                color=TRENCH, size_pt=12, bold=True,
            )
        self.add_heading_textbox(
            slide, title,
            left=text_left, top=1.0, width=text_width, height=1.4,
            color=TRENCH, size_pt=36, italic=True,
        )
        self._add_rich_paragraphs(
            slide,
            text=body,
            left=text_left, top=2.55, width=text_width, height=2.7,
            color=TRENCH, size_pt=14,
        )
        return slide

    def _paint_solid_panel(
        self,
        slide: Slide,
        left: float, top: float, width: float, height: float,
        color: RGBColor,
    ) -> None:
        """Solid brand-colour panel with a calm corner cluster.

        The brand's own Thank-You and Quote slides pair a solid panel with a single
        aquatic cluster tucked into a corner — not a full-bleed form that fills the
        whole panel and creates busy reflections. We follow that pattern here.
        The cluster is anchored to the bottom-right corner of the panel and sized
        to ~70% of the panel's shorter dimension so it reads as accent, not wallpaper.
        """
        block = self.add_rect(slide, left, top, width, height, fill=color)
        self._send_to_back(block)
        try:
            cluster_h = min(height, width) * 0.85
            cluster_w = cluster_h  # keep the aquatic cluster roughly square
            # Anchor to bottom-right of the panel, allowing a small bleed past the edge
            # so the form feels like it's crossing the panel boundary (brand move).
            clx = left + width - cluster_w * 0.75
            cly = top + height - cluster_h * 0.85
            pic = self.add_aquatic_form(
                slide, name="cluster-tall-01",
                left=clx, top=cly, width=cluster_w, height=cluster_h,
            )
            self._send_to_back(pic)
            self._send_to_back(block)  # keep block behind the form
        except FileNotFoundError:
            pass

    # ---- Team slides -------------------------------------------------------

    def add_team_member(
        self,
        name: str,
        role: str,
        bio: str,
        photo: Optional[str] = None,
        contact: Optional[str] = None,
        eyebrow: Optional[str] = "About our people",
    ) -> Slide:
        """
        Team member profile — portrait on the left, name + role + bio on the right.
        Modelled on brand slide-35.

        `photo`  : name from assets/portraits/ (e.g. "leader-1"). If missing, a branded
                   avatar placeholder is drawn instead.
        `contact`: optional contact line (e.g. "name@mearra.com, +358 ...") rendered small
                   at the bottom of the text column.
        """
        slide = self._add_slide("TITLE_ONLY")
        self._fill_placeholder(slide, 0, "")
        if eyebrow:
            self._fill_placeholder(slide, 2, eyebrow)

        # Portrait panel — full-bleed from the top edge to the bottom edge of the slide
        # and flush to the left edge, like brand slide-35. The portrait is the visual
        # anchor; the text column sits to its right on the Mist background.
        portrait_left = 0.0
        portrait_top = 0.0
        portrait_w = 4.20
        portrait_h = 5.625
        if photo:
            try:
                path = self._find_asset("portraits", photo, prefix="person-")
                pic = self.add_image(
                    slide, path, left=portrait_left, top=portrait_top,
                    width=portrait_w, height=portrait_h,
                    fit="cover",  # preserve aspect; centre-crop overflow
                )
            except FileNotFoundError:
                self._draw_avatar_placeholder(
                    slide, portrait_left + 0.5, 0.8, portrait_w - 1.0, portrait_h - 1.6,
                    initials=_initials(name),
                )
        else:
            self._draw_avatar_placeholder(
                slide, portrait_left + 0.5, 0.8, portrait_w - 1.0, portrait_h - 1.6,
                initials=_initials(name),
            )

        # Text column
        tx_left = portrait_w + 0.40
        tx_width = 10.0 - tx_left - 0.45
        self.add_body_textbox(
            slide, name,
            left=tx_left, top=0.9, width=tx_width, height=0.4,
            color=TRENCH, size_pt=14, bold=True,
        )
        self.add_heading_textbox(
            slide, role,
            left=tx_left, top=1.25, width=tx_width, height=0.7,
            color=TRENCH, size_pt=28, italic=True,
        )
        self._add_rich_paragraphs(
            slide, text=bio,
            left=tx_left, top=2.05, width=tx_width, height=2.8,
            color=TRENCH, size_pt=13,
        )
        if contact:
            self.add_body_textbox(
                slide, contact,
                left=tx_left, top=4.95, width=tx_width, height=0.35,
                color=TRENCH, size_pt=11,
            )
        return slide

    def add_team_grid(
        self,
        title: str,
        people: Iterable[Dict[str, Any]],
        eyebrow: Optional[str] = None,
    ) -> Slide:
        """
        Grid of circular team avatars (3 per row, up to 9). Each person dict:
            {"name": "Jean Duckling", "role": "Tech Support", "photo": "<optional portrait name>"}
        Modelled on brand slide-36.

        Missing photos fall back to a Trench-filled circle with the Mearra M — the same
        treatment the sample deck uses when a portrait isn't available.
        """
        people = list(people)
        if not 1 <= len(people) <= 9:
            raise ValueError("add_team_grid supports between 1 and 9 people")
        slide = self._add_slide("TITLE_ONLY")
        self._fill_placeholder(slide, 0, title)
        if eyebrow:
            self._fill_placeholder(slide, 2, eyebrow)

        # 3-column grid. Top row @ y=1.6, rows spaced 1.30" apart.
        avatar_size = 0.95
        col_lefts = [0.45, 3.55, 6.65]
        text_offset_x = avatar_size + 0.18
        row_top_0 = 1.60
        row_gap = 1.30

        for i, person in enumerate(people):
            row, col = divmod(i, 3)
            ax = col_lefts[col]
            ay = row_top_0 + row * row_gap
            # Avatar
            photo = person.get("photo")
            drew_photo = False
            if photo:
                try:
                    path = self._find_asset("portraits", photo, prefix="person-")
                    self._draw_circular_avatar_from_image(slide, path, ax, ay, avatar_size)
                    drew_photo = True
                except FileNotFoundError:
                    pass
            if not drew_photo:
                self._draw_avatar_placeholder(
                    slide, ax, ay, avatar_size, avatar_size,
                    initials=_initials(person.get("name", "")),
                    circular=True,
                    size_pt=34,  # brand spec: clearly readable but not dominant
                )
            # Name + role
            tx_left = ax + text_offset_x
            tx_width = (col_lefts[col + 1] - 0.1 - tx_left) if col < 2 else (10.0 - tx_left - 0.3)
            self.add_body_textbox(
                slide, person.get("name", ""),
                left=tx_left, top=ay + 0.08, width=tx_width, height=0.4,
                color=TRENCH, size_pt=14, bold=True,
            )
            self.add_body_textbox(
                slide, person.get("role", ""),
                left=tx_left, top=ay + 0.48, width=tx_width, height=0.4,
                color=TRENCH, size_pt=12,
            )
        return slide

    def _draw_avatar_placeholder(
        self,
        slide: Slide,
        left: float, top: float, width: float, height: float,
        initials: str = "",
        circular: bool = True,
        size_pt: Optional[int] = None,
    ) -> None:
        """Trench-filled oval/rect with the Mearra-style M (initials) — fallback for missing portraits.

        The initials sit in Sofia Sans Extra Condensed Black Italic (FONT_HEADING)
        and are centred both horizontally and vertically inside the circle. `size_pt`
        overrides the auto-sized font so callers can pin the glyph to a specific
        visual weight — e.g. the brand team grid uses 34pt inside a ~0.95" circle.
        """
        shape_kind = MSO_SHAPE.OVAL if circular else MSO_SHAPE.RECTANGLE
        shp = slide.shapes.add_shape(
            shape_kind, Inches(left), Inches(top), Inches(width), Inches(height),
        )
        shp.fill.solid()
        shp.fill.fore_color.rgb = TRENCH
        shp.line.fill.background()
        if initials:
            tf = shp.text_frame
            tf.margin_left = tf.margin_right = Emu(0)
            tf.margin_top = tf.margin_bottom = Emu(0)
            tf.word_wrap = False
            tf.vertical_anchor = MSO_ANCHOR.MIDDLE
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            r = p.add_run()
            r.text = initials
            r.font.name = FONT_HEADING
            r.font.bold = True
            r.font.italic = True
            if size_pt is None:
                # Auto-size: ~36pt per inch of diameter, clamped so tiny circles don't overflow.
                size_pt = int(max(12, min(64, height * 36)))
            r.font.size = Pt(size_pt)
            r.font.color.rgb = MIST

    def _draw_circular_avatar_from_image(
        self,
        slide: Slide,
        image_path: Path,
        left: float, top: float, size: float,
    ) -> None:
        """
        Draw a circular avatar by clipping an image into an oval shape using picture fill.
        Produces the rounded portrait treatment used in slide-35/-36.
        """
        shp = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Inches(left), Inches(top), Inches(size), Inches(size),
        )
        shp.line.fill.background()
        # Replace solid fill with picture fill via raw XML.
        spPr = shp.fill._xPr.spPr if hasattr(shp.fill, "_xPr") else shp.fill._xPr  # defensive
        # Simpler approach: drop the <a:solidFill> and inject <a:blipFill>.
        sp = shp._element
        spPr = sp.find(qn("p:spPr"))
        if spPr is None:
            return
        # remove existing fill elements
        for tag in ("a:noFill", "a:solidFill", "a:gradFill", "a:blipFill", "a:pattFill"):
            for el in spPr.findall(qn(tag)):
                spPr.remove(el)
        # insert blipFill referencing the image
        with open(image_path, "rb") as fh:
            image_bytes = fh.read()
        image_part, rId = slide.part.get_or_add_image(str(image_path))
        blipFill = etree.SubElement(spPr, qn("a:blipFill"))
        blip = etree.SubElement(blipFill, qn("a:blip"))
        blip.set(qn("r:embed"), rId)
        stretch = etree.SubElement(blipFill, qn("a:stretch"))
        etree.SubElement(stretch, qn("a:fillRect"))
        # Move blipFill to appear before <a:ln> if present.
        ln = spPr.find(qn("a:ln"))
        if ln is not None:
            spPr.remove(blipFill)
            spPr.insert(list(spPr).index(ln), blipFill)

    # ---- Timeline ----------------------------------------------------------

    def add_timeline(
        self,
        milestones: Iterable[Dict[str, str]],
        eyebrow: Optional[str] = None,
    ) -> Slide:
        """
        Brand **Timeline V1** — horizontal segmented bar with year labels above,
        colour-matched dots dropping below the segment boundaries, and numbered
        description blocks below the dots.

        milestones = [{"date": "2020", "body": "Lorem ipsum..."}, ...]

        * 3–6 milestones recommended (5 is the brand exemplar).
        * The bar is a sequence of solid-colour segments cycling through the brand
          palette: BLOSSOM → CORAL → INDIGO → CURRENT → LAGOON → TRENCH. The first
          segment is MIST (a "pre-history" stub), matching the brand slide.
        * Each year sits above a segment boundary; each dot drops below the same
          boundary with a thin TRENCH vertical lead.
        """
        milestones = list(milestones)
        if not 3 <= len(milestones) <= 6:
            raise ValueError("add_timeline supports 3–6 milestones")
        slide = self._add_slide("TITLE_ONLY")
        # Small eyebrow top-left only — no big title, per brand slide.
        self._fill_placeholder(slide, 0, "")
        self._fill_placeholder(slide, 2, eyebrow or "Timeline V1")

        # Segmented bar geometry — sits roughly mid-height, full-bleed width.
        bar_left = 0.0
        bar_right = 10.0
        bar_y = 2.30
        bar_h = 0.14
        n = len(milestones)
        # One "stub" segment at the very left + one segment per year.
        stub_w = 0.30
        segs_w = (bar_right - bar_left - stub_w) / n

        # Stub segment — MIST/grey, represents "before the first year."
        self.add_rect(slide, bar_left, bar_y, stub_w, bar_h, fill=MIST)

        # Year segments — cycle through brand palette. Order matches brand slide.
        seg_palette = [BLOSSOM, CORAL, INDIGO, CURRENT, LAGOON, TRENCH]
        seg_colors = [seg_palette[i % len(seg_palette)] for i in range(n)]

        # Year label baseline (above bar)
        year_baseline_y = bar_y - 0.85
        # Dot + number + body block below bar
        dot_d = 0.24
        drop_top = bar_y + bar_h
        drop_h = 0.65  # vertical lead from bar to dot
        num_y = drop_top + drop_h + 0.05
        body_y = num_y + 0.55

        block_w = max(1.60, segs_w - 0.15)  # give blocks a bit of breathing room

        for i, m in enumerate(milestones):
            # Segment
            seg_left = bar_left + stub_w + i * segs_w
            self.add_rect(slide, seg_left, bar_y, segs_w, bar_h, fill=seg_colors[i])

            # Anchor x: left edge of the segment (where the label and dot align)
            anchor_x = seg_left

            # Year label above the anchor
            self.add_body_textbox(
                slide, m.get("date", ""),
                left=anchor_x - 0.35, top=year_baseline_y,
                width=1.80, height=0.55,
                color=TRENCH, size_pt=22, bold=True,
            )

            # Vertical lead (thin line) dropping from the segment boundary to the dot
            lead = slide.shapes.add_connector(
                1,  # STRAIGHT
                Inches(anchor_x), Inches(drop_top),
                Inches(anchor_x), Inches(drop_top + drop_h),
            )
            lead.line.color.rgb = seg_colors[i]
            lead.line.width = Pt(1.0)

            # Dot below the lead
            dot = slide.shapes.add_shape(
                MSO_SHAPE.OVAL,
                Inches(anchor_x - dot_d / 2), Inches(drop_top + drop_h - dot_d / 2),
                Inches(dot_d), Inches(dot_d),
            )
            dot.fill.solid()
            dot.fill.fore_color.rgb = seg_colors[i]
            dot.line.fill.background()

            # Number (01, 02, ...)
            self.add_body_textbox(
                slide, f"{i + 1:02d}",
                left=anchor_x - 0.10, top=num_y,
                width=1.00, height=0.45,
                color=TRENCH, size_pt=20, bold=True,
            )

            # Body text
            body = m.get("body", "")
            if body:
                self.add_body_textbox(
                    slide, body,
                    left=anchor_x - 0.10, top=body_y,
                    width=block_w, height=2.0,
                    color=TRENCH, size_pt=13,
                )
        return slide

    def add_timeline_arrows(
        self,
        title: str,
        milestones: Iterable[Dict[str, str]],
        eyebrow: Optional[str] = None,
    ) -> Slide:
        """
        Brand **Timeline V2** — chevron arrow track. Each milestone is rendered as a
        right-pointing chevron arrow in an indigo-tint gradient (pale → saturated, left
        to right). Year labels hang from TRENCH lollipop pins *above* each chevron,
        and a numbered description block sits below.

        milestones = [{"date": "2023", "body": "Lorem ipsum..."}, ...]

        * 3–5 milestones; 4 reads cleanest.
        * The chevron strip runs at roughly 40% of slide height so the title above
          and the numbered blocks below both have comfortable breathing room.
        """
        milestones = list(milestones)
        if not 3 <= len(milestones) <= 5:
            raise ValueError("add_timeline_arrows supports 3–5 milestones")
        slide = self._add_slide("TITLE_ONLY")
        self._fill_placeholder(slide, 0, title)
        if eyebrow:
            self._fill_placeholder(slide, 2, eyebrow)

        # Chevron strip geometry — centred across the slide, leaving margin left/right.
        n = len(milestones)
        strip_left = 1.40
        strip_right = 9.60
        strip_w = strip_right - strip_left
        chev_h = 0.55
        chev_y = 2.45
        chev_w = strip_w / n

        tints = _indigo_tints(n)  # pale → saturated

        # Year label + lollipop pins above each chevron.
        pin_top = 1.10
        pin_len = chev_y - pin_top - 0.05  # line length
        num_y = chev_y + chev_h + 0.40
        body_y = num_y + 0.55

        for i, m in enumerate(milestones):
            cx = strip_left + i * chev_w
            # Chevron shape — CHEVRON points right.
            ch = slide.shapes.add_shape(
                MSO_SHAPE.CHEVRON,
                Inches(cx), Inches(chev_y),
                Inches(chev_w), Inches(chev_h),
            )
            ch.fill.solid()
            ch.fill.fore_color.rgb = tints[i]
            ch.line.fill.background()

            # Anchor for year label: horizontal middle of the chevron body.
            anchor_x = cx + chev_w * 0.45

            # Lollipop pin: small TRENCH dot + vertical line down to the chevron top.
            pin_d = 0.12
            pin_dot = slide.shapes.add_shape(
                MSO_SHAPE.OVAL,
                Inches(anchor_x - pin_d / 2), Inches(pin_top + 0.45 - pin_d / 2),
                Inches(pin_d), Inches(pin_d),
            )
            pin_dot.fill.solid()
            pin_dot.fill.fore_color.rgb = TRENCH
            pin_dot.line.fill.background()
            pin_line = slide.shapes.add_connector(
                1,
                Inches(anchor_x), Inches(pin_top + 0.45 + pin_d / 2),
                Inches(anchor_x), Inches(chev_y),
            )
            pin_line.line.color.rgb = TRENCH
            pin_line.line.width = Pt(1.0)

            # Year label above the pin
            self.add_body_textbox(
                slide, m.get("date", ""),
                left=anchor_x - 0.70, top=pin_top,
                width=1.40, height=0.45,
                color=TRENCH, size_pt=18, bold=True,
            )

            # Number + body below the chevron
            block_left = cx - 0.10
            block_w = chev_w + 0.05
            self.add_body_textbox(
                slide, f"{i + 1:02d}",
                left=block_left, top=num_y,
                width=1.00, height=0.45,
                color=TRENCH, size_pt=20, bold=True,
            )
            body = m.get("body", "")
            if body:
                self.add_body_textbox(
                    slide, body,
                    left=block_left, top=body_y,
                    width=block_w, height=2.0,
                    color=TRENCH, size_pt=12,
                )
        return slide

    # ---- Table -------------------------------------------------------------

    def add_table(
        self,
        title: str,
        headers: Iterable[str],
        rows: Iterable[Iterable[str]],
        eyebrow: Optional[str] = None,
        col_widths: Optional[Iterable[float]] = None,
    ) -> Slide:
        """
        Title + data table. Header row sits on Indigo with Mist text; data rows alternate
        White and Mist for readable zebra striping. Sofia Sans throughout.

        `col_widths` : optional list of inches per column. If omitted, columns are even.
        """
        headers = list(headers)
        rows = [list(r) for r in rows]
        n_cols = len(headers)
        n_rows = 1 + len(rows)
        if n_cols == 0:
            raise ValueError("Table needs at least one column")

        slide = self._add_slide("TITLE_ONLY")
        self._fill_placeholder(slide, 0, title)
        if eyebrow:
            self._fill_placeholder(slide, 2, eyebrow)

        # Table geometry
        tbl_left = 0.55
        tbl_top = 1.80
        tbl_width = 8.90
        tbl_height = min(3.40, 0.50 + len(rows) * 0.42)
        shp = slide.shapes.add_table(
            rows=n_rows, cols=n_cols,
            left=Inches(tbl_left), top=Inches(tbl_top),
            width=Inches(tbl_width), height=Inches(tbl_height),
        )
        table = shp.table

        # Column widths
        if col_widths:
            col_widths = list(col_widths)
            if len(col_widths) != n_cols:
                raise ValueError("col_widths must match number of columns")
            for i, w in enumerate(col_widths):
                table.columns[i].width = Inches(w)

        # Row heights — header a bit taller than body
        table.rows[0].height = Inches(0.50)
        for r in range(1, n_rows):
            table.rows[r].height = Inches(0.38)

        def style_cell(cell, text, *, bold=False, color=TRENCH, fill=WHITE, size_pt=12):
            cell.fill.solid()
            cell.fill.fore_color.rgb = fill
            cell.margin_left = Inches(0.14)
            cell.margin_right = Inches(0.14)
            cell.margin_top = Inches(0.06)
            cell.margin_bottom = Inches(0.06)
            tf = cell.text_frame
            tf.word_wrap = True
            tf.vertical_anchor = MSO_ANCHOR.MIDDLE
            p = tf.paragraphs[0]
            p.text = ""  # clear default
            r = p.add_run()
            r.text = text
            r.font.name = FONT_BODY
            r.font.bold = bold
            r.font.size = Pt(size_pt)
            r.font.color.rgb = color

        # Header row
        for c, h in enumerate(headers):
            style_cell(table.cell(0, c), h, bold=True, color=MIST, fill=INDIGO, size_pt=13)

        # Data rows (zebra)
        for r, row in enumerate(rows, start=1):
            fill = WHITE if r % 2 == 1 else MIST
            for c in range(n_cols):
                value = row[c] if c < len(row) else ""
                style_cell(table.cell(r, c), str(value), fill=fill, color=TRENCH)

        return slide

    # ---- Chart-style slides -----------------------------------------------

    def add_pie_chart(
        self,
        title: str,
        slices: Iterable[Dict[str, Any]],
        eyebrow: Optional[str] = None,
        donut: bool = True,
        legend: bool = True,
    ) -> Slide:
        """
        Title + pie / donut chart on the right, brand legend on the left.

        slices = [{"label": "Development", "value": 55, "color": CORAL}, ...]
        Colours default to a cycle through the brand palette (Indigo, Coral, Blossom,
        Current, Lagoon, Trench).

        Implementation uses python-pptx's native chart support so the output renders
        correctly in PowerPoint, Keynote, Google Slides, and LibreOffice.
        """
        slices = list(slices)
        if not 2 <= len(slices) <= 8:
            raise ValueError("add_pie_chart supports 2–8 slices")
        slide = self._add_slide("TITLE_ONLY")
        self._fill_placeholder(slide, 0, title)
        if eyebrow:
            self._fill_placeholder(slide, 2, eyebrow)

        palette = [INDIGO, CORAL, BLOSSOM, CURRENT, LAGOON, TRENCH, LAGOON, BLOSSOM]
        for i, s in enumerate(slices):
            s.setdefault("color", palette[i % len(palette)])
        total = sum(float(s.get("value", 0)) for s in slices) or 1.0

        chart_data = CategoryChartData()
        chart_data.categories = [s.get("label", "") for s in slices]
        chart_data.add_series("Slices", [float(s.get("value", 0)) for s in slices])

        chart_type = XL_CHART_TYPE.DOUGHNUT if donut else XL_CHART_TYPE.PIE
        # Place the chart on the right half of the slide.
        cx, cy, cw, ch = 5.0, 1.40, 4.80, 3.80
        gframe = slide.shapes.add_chart(
            chart_type, Inches(cx), Inches(cy), Inches(cw), Inches(ch),
            chart_data,
        )
        chart = gframe.chart
        chart.has_title = False
        chart.has_legend = False  # we draw our own brand legend on the left

        # Color each slice (data point) with the brand palette.
        plot = chart.plots[0]
        series = plot.series[0]
        for i, s in enumerate(slices):
            try:
                point = series.points[i]
                fill = point.format.fill
                fill.solid()
                fill.fore_color.rgb = s["color"]
                point.format.line.color.rgb = MIST
                point.format.line.width = Pt(1.5)
            except Exception:
                # Some python-pptx versions don't expose data_labels per point — skip.
                pass

        # Brand legend on the left — colour swatch + label + percent.
        if legend:
            leg_left = 0.55
            leg_top = 1.80
            row_h = 0.44
            for i, s in enumerate(slices):
                y = leg_top + i * row_h
                self.add_rect(
                    slide, left=leg_left, top=y + 0.08, width=0.35, height=0.22,
                    fill=s["color"],
                )
                pct = 100.0 * float(s.get("value", 0)) / total
                self.add_body_textbox(
                    slide, s.get("label", ""),
                    left=leg_left + 0.50, top=y, width=3.0, height=0.40,
                    color=TRENCH, size_pt=13, bold=True,
                )
                self.add_body_textbox(
                    slide, f"{pct:.0f}%",
                    left=leg_left + 3.55, top=y, width=0.80, height=0.40,
                    color=TRENCH, size_pt=13, bold=True,
                )
        return slide

    def _draw_pie_slice(
        self,
        slide,
        cx: float, cy: float,
        r_outer: float, r_inner: float,
        start_deg: float, sweep_deg: float,
        fill: RGBColor,
    ):
        """
        Draw one pie/donut slice as a freeform shape using EMUs. Angles in degrees,
        cx/cy/r in inches. start_deg=-90 places 0 at 12 o'clock. Positive sweep = clockwise.
        """
        emu_per_in = 914400

        def pt(angle_deg: float, radius: float):
            rad = math.radians(angle_deg)
            return (cx + radius * math.cos(rad), cy + radius * math.sin(rad))

        end_deg = start_deg + sweep_deg
        steps = max(6, int(abs(sweep_deg) / 5))
        outer_pts = [pt(start_deg + (end_deg - start_deg) * i / steps, r_outer)
                     for i in range(steps + 1)]
        if r_inner > 0:
            inner_pts = [pt(end_deg - (end_deg - start_deg) * i / steps, r_inner)
                         for i in range(steps + 1)]
            all_pts = outer_pts + inner_pts
        else:
            all_pts = outer_pts + [(cx, cy)]

        # build_freeform: start_x/y and subsequent points are multiplied by `scale`
        # to get EMU. Pass raw inch values (not Inches()), scale=EMU/inch.
        first = all_pts[0]
        remaining = [(p[0], p[1]) for p in all_pts[1:]]
        freeform = slide.shapes.build_freeform(
            first[0], first[1], scale=emu_per_in,
        )
        freeform.add_line_segments(remaining, close=True)
        shp = freeform.convert_to_shape()
        shp.fill.solid()
        shp.fill.fore_color.rgb = fill
        shp.line.color.rgb = MIST
        shp.line.width = Pt(1.5)
        return shp

    def add_pyramid(
        self,
        title: str,
        levels: Iterable[Dict[str, str]],
        eyebrow: Optional[str] = None,
        inverted: bool = False,
    ) -> Slide:
        """
        3–5 stacked pyramid levels with labels to the right of each tier. Bottom tier is
        widest; each tier above shrinks by a consistent step. Labels are Sofia Sans.

        levels = [{"label": "Strategy", "body": "..."}, ...]
        Order in the list = top → bottom. Set `inverted=True` to flip (bottom → top).
        Bottom tier uses Indigo, top tier uses Blossom, middle tiers interpolate through
        Coral / Current — purely brand palette.
        """
        levels = list(levels)
        if not 3 <= len(levels) <= 5:
            raise ValueError("add_pyramid supports 3–5 levels")
        if inverted:
            levels = list(reversed(levels))
        n = len(levels)
        slide = self._add_slide("TITLE_ONLY")
        self._fill_placeholder(slide, 0, title)
        if eyebrow:
            self._fill_placeholder(slide, 2, eyebrow)

        # Pyramid geometry: left-half of slide, centred horizontally
        cx = 2.85
        top_y = 1.7
        total_h = 3.5
        tier_h = total_h / n
        min_w = 0.9   # top tier width
        max_w = 4.2   # bottom tier width
        # Tier colours — pick n colours from a brand-safe gradient.
        tier_palette = _sample_brand_gradient(n)

        for i, lvl in enumerate(levels):
            # i=0 is top (smallest), i=n-1 is bottom (largest)
            w_top = min_w + (max_w - min_w) * (i / max(1, n))
            w_bot = min_w + (max_w - min_w) * ((i + 1) / max(1, n))
            # Use TRAPEZOID shape, positioned so it widens downward.
            # MSO_SHAPE.TRAPEZOID has a narrower top, wider bottom by default.
            left = cx - w_bot / 2
            top = top_y + i * tier_h
            # We size the bounding box to the wider bottom; the trapezoid widens to the full
            # bottom width. The top edge of the trapezoid is at ~75% of the box width by
            # default — we compensate by sizing the shape to the bottom width and accept
            # the visual step between tiers.
            shp = slide.shapes.add_shape(
                MSO_SHAPE.TRAPEZOID,
                Inches(left), Inches(top), Inches(w_bot), Inches(tier_h),
            )
            # Rotate 180° so the wider edge faces up — no, that flips it; we actually want
            # a *pyramid*, which means top narrow, bottom wide — that's the default
            # TRAPEZOID orientation. Leave as is.
            shp.fill.solid()
            shp.fill.fore_color.rgb = tier_palette[i]
            shp.line.fill.background()
            # Tier label inside the tier
            tf = shp.text_frame
            tf.margin_left = tf.margin_right = Inches(0.1)
            tf.margin_top = tf.margin_bottom = Inches(0.03)
            tf.vertical_anchor = MSO_ANCHOR.MIDDLE
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            r = p.add_run()
            r.text = lvl.get("label", "")
            r.font.name = FONT_HEADING
            r.font.bold = True
            r.font.italic = True
            r.font.size = Pt(max(14, int(22 - (n - 3) * 2)))
            # Text colour — light on dark tiers, dark on light tiers.
            r.font.color.rgb = MIST if tier_palette[i] in (INDIGO, TRENCH, CORAL) else TRENCH

            # Body caption to the right of the pyramid (aligned with this tier).
            # Accept either "body" or "description" so callers can use whichever
            # feels natural for their content.
            caption = lvl.get("body") or lvl.get("description")
            if caption:
                self.add_body_textbox(
                    slide, caption,
                    left=5.75, top=top + 0.05, width=3.85, height=tier_h - 0.1,
                    color=TRENCH, size_pt=12,
                )
        return slide

    def add_circular_lifecycle(
        self,
        phases: Iterable[Dict[str, Any]],
        center_label: Optional[str] = None,
        title: Optional[str] = None,
        eyebrow: Optional[str] = None,
    ) -> Slide:
        """
        The brand "Proactive strategic digital partner" service-offering slide.

        A centred indigo-gradient donut with four numbered quadrant labels, each
        connected to the ring by a short lead line. The centre of the donut carries
        a bold uppercase multi-line label. Below each numbered title sits a short
        list of sub-items (the services inside that quadrant).

        phases = [
            {"title": "Research & Discovery", "items": ["Audits", "Prototypes", ...]},
            {"title": "Design & Delivery",    "items": [...]},
            {"title": "Improve & Optimize",   "items": [...]},
            {"title": "Maintain & Secure",    "items": [...]},
        ]

        * `phases` are listed in the drawing order 1→4: top-right, bottom-right,
          bottom-left, top-left. This matches the brand slide's reading order.
        * `center_label` is rendered uppercase, bold, in TRENCH inside the donut hole.
        * `title` / `eyebrow` are optional; the brand slide leaves them off by default
          because the centre label carries the message.
        """
        phases = list(phases)
        if len(phases) != 4:
            raise ValueError("add_circular_lifecycle requires exactly 4 phases")
        slide = self._add_slide("TITLE_ONLY")
        self._fill_placeholder(slide, 0, title or "")
        self._fill_placeholder(slide, 2, eyebrow or "")

        # ---- Ring geometry (centred) -------------------------------------
        slide_w, slide_h = 10.0, 5.625
        ring_d = 3.40
        ring_left = (slide_w - ring_d) / 2
        ring_top = (slide_h - ring_d) / 2
        center_x = slide_w / 2
        center_y = slide_h / 2

        # Outer ring — solid INDIGO with a soft radial gradient overlay for depth.
        # python-pptx doesn't expose radial gradients cleanly, so we fake the
        # effect: draw a saturated INDIGO outer oval, then overlay a pale tint oval
        # along the top-right to suggest the brand's soft gradient highlight.
        outer = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Inches(ring_left), Inches(ring_top),
            Inches(ring_d), Inches(ring_d),
        )
        outer.fill.solid()
        outer.fill.fore_color.rgb = INDIGO
        outer.line.fill.background()
        # Apply a gradient fill via raw XML for a more brand-faithful look.
        self._apply_indigo_gradient(outer)

        # Inner hole — MIST oval cut out of the centre.
        hole_d = 1.55
        hole_left = center_x - hole_d / 2
        hole_top = center_y - hole_d / 2
        hole = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Inches(hole_left), Inches(hole_top),
            Inches(hole_d), Inches(hole_d),
        )
        hole.fill.solid()
        hole.fill.fore_color.rgb = MIST
        hole.line.fill.background()

        # Thin TRENCH cross-hairs to suggest the quadrant boundaries — they
        # appear as tiny dark ticks where they intersect the ring edge.
        # Vertical
        vline = slide.shapes.add_connector(
            1,
            Inches(center_x), Inches(ring_top),
            Inches(center_x), Inches(ring_top + ring_d),
        )
        vline.line.color.rgb = MIST
        vline.line.width = Pt(1.5)
        # Horizontal
        hline = slide.shapes.add_connector(
            1,
            Inches(ring_left), Inches(center_y),
            Inches(ring_left + ring_d), Inches(center_y),
        )
        hline.line.color.rgb = MIST
        hline.line.width = Pt(1.5)

        # Centre label
        if center_label:
            self.add_body_textbox(
                slide, center_label.upper(),
                left=hole_left + 0.05, top=hole_top + 0.20,
                width=hole_d - 0.10, height=hole_d - 0.40,
                color=TRENCH, size_pt=11, bold=True,
            )

        # ---- Quadrant labels + lead lines --------------------------------
        # Each quadrant: (label_left, label_top, label_width, align, dot_dx, dot_dy)
        # where (dot_dx, dot_dy) is the lead-line anchor point on the ring edge
        # as an offset from the ring centre (radius 1 unit, so the unit circle).
        # Order: top-right, bottom-right, bottom-left, top-left.
        sqrt2 = math.sqrt(2) / 2
        # Ring spans x = 3.30 → 6.70 (ring_left → ring_left + ring_d).
        # Labels must start ≥ 6.80 on the right and end ≤ 3.20 on the left to clear
        # the ring. Keep a 0.10" breathing gap so the lead lines are visible.
        quadrants = [
            # 1. Top-right — label sits in upper-right, text left-aligned
            {"lx": 6.85, "ly": 1.25, "lw": 3.00, "align": "left",
             "unit_dx":  sqrt2, "unit_dy": -sqrt2},
            # 2. Bottom-right
            {"lx": 6.85, "ly": 3.20, "lw": 3.00, "align": "left",
             "unit_dx":  sqrt2, "unit_dy":  sqrt2},
            # 3. Bottom-left — text right-aligned
            {"lx": 0.15, "ly": 3.20, "lw": 3.00, "align": "right",
             "unit_dx": -sqrt2, "unit_dy":  sqrt2},
            # 4. Top-left — text right-aligned
            {"lx": 0.15, "ly": 1.25, "lw": 3.00, "align": "right",
             "unit_dx": -sqrt2, "unit_dy": -sqrt2},
        ]

        for i, (q, phase) in enumerate(zip(quadrants, phases)):
            align = PP_ALIGN.LEFT if q["align"] == "left" else PP_ALIGN.RIGHT

            # Lead line: from ring-edge dot to the inner edge of the label.
            r = ring_d / 2
            ring_edge_x = center_x + q["unit_dx"] * r
            ring_edge_y = center_y + q["unit_dy"] * r
            # Label inner anchor (near edge of label text, close to ring)
            if q["align"] == "left":
                label_anchor_x = q["lx"]
            else:
                label_anchor_x = q["lx"] + q["lw"]
            label_anchor_y = q["ly"] + 0.18  # near top of title

            # Short TRENCH lead line from label anchor → ring edge
            lead = slide.shapes.add_connector(
                1,
                Inches(label_anchor_x), Inches(label_anchor_y),
                Inches(ring_edge_x), Inches(ring_edge_y),
            )
            lead.line.color.rgb = TRENCH
            lead.line.width = Pt(1.0)

            # TRENCH dot at the ring edge
            dot_d = 0.11
            dot = slide.shapes.add_shape(
                MSO_SHAPE.OVAL,
                Inches(ring_edge_x - dot_d / 2), Inches(ring_edge_y - dot_d / 2),
                Inches(dot_d), Inches(dot_d),
            )
            dot.fill.solid()
            dot.fill.fore_color.rgb = TRENCH
            dot.line.fill.background()

            # Numbered title — bold TRENCH (per brand; NOT indigo)
            title_tb = slide.shapes.add_textbox(
                Inches(q["lx"]), Inches(q["ly"]),
                Inches(q["lw"]), Inches(0.45),
            )
            tf = title_tb.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.alignment = align
            r1 = p.add_run()
            r1.text = f"{i + 1}. {phase.get('title', '')}"
            r1.font.name = FONT_BODY
            r1.font.bold = True
            r1.font.size = Pt(15)
            r1.font.color.rgb = TRENCH

            # Items list — TRENCH regular
            items = phase.get("items", []) or []
            if items:
                items_tb = slide.shapes.add_textbox(
                    Inches(q["lx"]), Inches(q["ly"] + 0.50),
                    Inches(q["lw"]), Inches(1.80),
                )
                tf = items_tb.text_frame
                tf.word_wrap = True
                for j, item in enumerate(items):
                    p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
                    p.alignment = align
                    r = p.add_run()
                    r.text = item
                    r.font.name = FONT_BODY
                    r.font.size = Pt(11)
                    r.font.color.rgb = TRENCH
                    p.space_after = Pt(2)
        return slide

    def _apply_indigo_gradient(self, shape) -> None:
        """
        Replace a shape's solid fill with a brand indigo linear gradient:
        pale indigo top-left → saturated INDIGO middle → soft pink at bottom-right.
        Matches the brand service-offering ring (slide-41).
        """
        sp = shape._element
        spPr = sp.find(qn("p:spPr"))
        if spPr is None:
            return
        # Remove any existing fill children
        for tag in ("a:noFill", "a:solidFill", "a:gradFill", "a:blipFill", "a:pattFill"):
            for el in spPr.findall(qn(tag)):
                spPr.remove(el)
        # Build <a:gradFill> with three stops.
        grad = etree.SubElement(spPr, qn("a:gradFill"))
        grad.set("flip", "none")
        grad.set("rotWithShape", "1")
        gsLst = etree.SubElement(grad, qn("a:gsLst"))
        stops = [
            ("0",     "C8BEFF"),   # pale violet top-left
            ("55000", "4400F0"),   # full INDIGO
            ("100000","F0B5D1"),   # faint blossom pink at bottom-right
        ]
        for pos, hex_ in stops:
            gs = etree.SubElement(gsLst, qn("a:gs"), pos=pos)
            etree.SubElement(gs, qn("a:srgbClr"), val=hex_)
        # 135° diagonal: 8100000 EMU angle units ≡ 135°
        lin = etree.SubElement(grad, qn("a:lin"))
        lin.set("ang", "8100000")
        lin.set("scaled", "1")
        # Re-insert the <a:ln> element after gradient if it exists
        ln = spPr.find(qn("a:ln"))
        if ln is not None:
            spPr.remove(ln)
            spPr.append(ln)

    # ---- Chart & text (split panel, brand slides 28–31) -------------------

    def add_chart_panel(
        self,
        title: str,
        visual: str,
        items: Iterable[Dict[str, Any]],
        subtitle: Optional[str] = None,
        intro: Optional[str] = None,
        center_label: Optional[str] = None,
        show_values: bool = False,
    ) -> Slide:
        """
        The brand's "Chart & text" split-panel slide — modelled exactly on brand
        slides 28 (pyramid), 29 (nested circles), 30 (cycle), and 31 (numbered donut).

        Layout: left half is a TRENCH-filled panel with the visual; right half is the
        default Mist background with a title, an INDIGO subtitle, an intro line, and
        a numbered list of the items in the visual. Numbers on the visual match the
        numbers in the list so the eye can connect them.

        visual   : "pyramid" | "circles" | "cycle" | "donut"
        items    : [{"label": "Eat", "value": 21.8}, ...]
                   `value` is required for "circles" and optional for "donut" — if
                   omitted, donut slices are equal. 3–5 items for pyramid, 4 for
                   circles/cycle/donut reads cleanest.
        subtitle : small INDIGO caption under the title, e.g. "Chart & text".
        intro    : line above the numbered list, e.g. "Daily duck activities:".
        center_label : for "cycle" — text inside the ring.
        show_values  : for "donut" — append "NN.N%" to each list item.
        """
        items = list(items)
        if visual not in {"pyramid", "circles", "cycle", "donut"}:
            raise ValueError("visual must be one of: pyramid, circles, cycle, donut")
        slide = self._add_slide("TITLE_ONLY")
        # Clear the master title; we draw our own in the right column.
        self._fill_placeholder(slide, 0, "")
        self._fill_placeholder(slide, 2, "")

        # ---- Left panel: Trench background -----------------------------------
        panel_left, panel_top, panel_w, panel_h = 0.0, 0.0, 4.50, 5.625
        panel = self.add_rect(slide, panel_left, panel_top, panel_w, panel_h, fill=TRENCH)
        self._send_to_back(panel)

        # ---- Right column text ----------------------------------------------
        tx_left = 4.95
        tx_width = 10.0 - tx_left - 0.45
        # Title
        self.add_body_textbox(
            slide, title,
            left=tx_left, top=0.85, width=tx_width, height=1.20,
            color=TRENCH, size_pt=26, bold=True,
        )
        y = 2.00
        if subtitle:
            self.add_body_textbox(
                slide, subtitle,
                left=tx_left, top=y, width=tx_width, height=0.40,
                color=INDIGO, size_pt=15, bold=True,
            )
            y += 0.55
        if intro:
            self.add_body_textbox(
                slide, intro,
                left=tx_left, top=y, width=tx_width, height=0.35,
                color=TRENCH, size_pt=13,
            )
            y += 0.45

        # Numbered list — mirrors the numbers on the visual.
        list_tb = slide.shapes.add_textbox(
            Inches(tx_left + 0.15), Inches(y), Inches(tx_width - 0.15), Inches(3.0),
        )
        tf = list_tb.text_frame
        tf.word_wrap = True
        for i, it in enumerate(items):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.alignment = PP_ALIGN.LEFT
            label = it.get("label", "")
            if show_values and it.get("value") is not None:
                label = f"{label} {float(it['value']):.1f}%"
            r = p.add_run()
            r.text = f"{i + 1}.  {label}"
            r.font.name = FONT_BODY
            r.font.size = Pt(13)
            r.font.color.rgb = TRENCH
            p.space_after = Pt(4)

        # ---- Visual inside the Trench panel ----------------------------------
        if visual == "pyramid":
            self._draw_chart_pyramid(slide, items, panel_left, panel_top, panel_w, panel_h)
        elif visual == "circles":
            self._draw_chart_nested_circles(slide, items, panel_left, panel_top, panel_w, panel_h)
        elif visual == "cycle":
            self._draw_chart_cycle_ring(
                slide, items, panel_left, panel_top, panel_w, panel_h,
                center_label=center_label,
            )
        elif visual == "donut":
            self._draw_chart_numbered_donut(slide, items, panel_left, panel_top, panel_w, panel_h)
        return slide

    # ---- Chart panel visuals (drawn inside a Trench left panel) -----------

    def _draw_chart_pyramid(
        self,
        slide: Slide,
        items: List[Dict[str, Any]],
        px: float, py: float, pw: float, ph: float,
    ) -> None:
        """Stacked trapezoid pyramid with numbered labels, inside a Trench panel.
        Layers get progressively darker indigo tints from top to bottom (brand slide-28).
        """
        n = len(items)
        if not 3 <= n <= 6:
            raise ValueError("chart_panel pyramid expects 3–6 items")
        # Centered on the left panel
        apex_x = px + pw / 2
        total_h = ph * 0.66              # use ~2/3 of panel height
        top_y = py + (ph - total_h) / 2
        min_w = pw * 0.12                # narrow apex
        max_w = pw * 0.76                # wide base
        tier_h = total_h / n

        # Colour ramp: lightest at top -> INDIGO at bottom. Use tinted versions of INDIGO.
        tints = _indigo_tints(n)  # list of RGBColor, pale → saturated

        for i, it in enumerate(items):
            w_bot = min_w + (max_w - min_w) * ((i + 1) / n)
            left = apex_x - w_bot / 2
            top = top_y + i * tier_h
            shp = slide.shapes.add_shape(
                MSO_SHAPE.TRAPEZOID,
                Inches(left), Inches(top), Inches(w_bot), Inches(tier_h),
            )
            shp.fill.solid()
            shp.fill.fore_color.rgb = tints[i]
            shp.line.fill.background()
            # Number label inside the tier — white on dark, trench on light.
            tf = shp.text_frame
            tf.vertical_anchor = MSO_ANCHOR.MIDDLE
            tf.margin_left = tf.margin_right = Inches(0.05)
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            r = p.add_run()
            r.text = str(i + 1)
            r.font.name = FONT_BODY
            r.font.bold = True
            r.font.size = Pt(16)
            # Light tints: use TRENCH; dark tints: use WHITE/MIST.
            r.font.color.rgb = MIST if _is_dark(tints[i]) else TRENCH

    def _draw_chart_nested_circles(
        self,
        slide: Slide,
        items: List[Dict[str, Any]],
        px: float, py: float, pw: float, ph: float,
    ) -> None:
        """Concentric nested circles anchored at the bottom of the left panel,
        each sized proportional to its value. Mirrors brand slide-29.
        """
        n = len(items)
        if not 2 <= n <= 5:
            raise ValueError("chart_panel circles expects 2–5 items")
        # Circles share a common baseline (bottom) so smaller ones "sit inside" bigger ones.
        max_d = min(pw * 0.76, ph * 0.82)
        baseline_y = py + ph - (ph - max_d) / 2  # bottom of the largest circle
        cx = px + pw / 2

        # Ordered from largest (first) to smallest (last) in the visual. Items keep their
        # natural order; the first item gets the biggest ring.
        tints = _indigo_tints(n)
        tints.reverse()  # biggest (outer) = pale, smallest (inner) = saturated
        for i, it in enumerate(items):
            # Value-proportional diameter; fallback to equal steps if no values.
            if it.get("value") is not None:
                v = float(it["value"])
                # Normalize so first item (max) fills max_d.
                max_v = max(float(x.get("value", 1)) for x in items) or 1.0
                d = max_d * (v / max_v)
            else:
                d = max_d * (1 - i / n)
            left = cx - d / 2
            top = baseline_y - d
            oval = slide.shapes.add_shape(
                MSO_SHAPE.OVAL,
                Inches(left), Inches(top), Inches(d), Inches(d),
            )
            oval.fill.solid()
            oval.fill.fore_color.rgb = tints[i]
            oval.line.fill.background()
            # Label inside the top of the circle — centered horizontally on the ring.
            label_text = it.get("label", "")
            if it.get("value") is not None:
                label_text = f"{float(it['value']):.0f}%"
            tb = slide.shapes.add_textbox(
                Inches(left), Inches(top + d * 0.05),
                Inches(d), Inches(0.55),
            )
            tf = tb.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            r = p.add_run()
            r.text = label_text
            r.font.name = FONT_BODY
            r.font.bold = True
            r.font.size = Pt(18)
            r.font.color.rgb = MIST if _is_dark(tints[i]) else TRENCH

    def _draw_chart_cycle_ring(
        self,
        slide: Slide,
        items: List[Dict[str, Any]],
        px: float, py: float, pw: float, ph: float,
        center_label: Optional[str] = None,
    ) -> None:
        """Ring with N equal segments in tinted indigo, center label overlay.
        Modelled on brand slide-30's 'A day in the life' cycle.
        """
        n = len(items)
        if not 2 <= n <= 8:
            raise ValueError("chart_panel cycle expects 2–8 items")
        # Centered ring
        ring_d = min(pw * 0.82, ph * 0.78)
        cx = px + (pw - ring_d) / 2
        cy = py + (ph - ring_d) / 2
        chart_data = CategoryChartData()
        chart_data.categories = [str(i + 1) for i in range(n)]
        chart_data.add_series("Cycle", [1] * n)
        gframe = slide.shapes.add_chart(
            XL_CHART_TYPE.DOUGHNUT,
            Inches(cx), Inches(cy), Inches(ring_d), Inches(ring_d),
            chart_data,
        )
        chart = gframe.chart
        chart.has_title = False
        chart.has_legend = False
        tints = _indigo_tints(n)
        try:
            series = chart.plots[0].series[0]
            for i in range(n):
                p = series.points[i]
                p.format.fill.solid()
                p.format.fill.fore_color.rgb = tints[i]
                p.format.line.color.rgb = TRENCH
                p.format.line.width = Pt(1.5)
        except Exception:
            pass

        if center_label:
            hole_w = ring_d * 0.40
            hole_h = 0.90
            self.add_heading_textbox(
                slide, center_label,
                left=cx + (ring_d - hole_w) / 2,
                top=cy + (ring_d - hole_h) / 2,
                width=hole_w, height=hole_h,
                color=MIST, size_pt=11, italic=False,
            )

    def _draw_chart_numbered_donut(
        self,
        slide: Slide,
        items: List[Dict[str, Any]],
        px: float, py: float, pw: float, ph: float,
    ) -> None:
        """Doughnut with numbered wedges in indigo tints.
        Numbers are drawn on the wedges at their angular midpoints, matching slide-31.
        """
        n = len(items)
        if not 2 <= n <= 8:
            raise ValueError("chart_panel donut expects 2–8 items")
        ring_d = min(pw * 0.82, ph * 0.78)
        cx = px + (pw - ring_d) / 2
        cy = py + (ph - ring_d) / 2
        # Values drive wedge sizes; default to equal if missing.
        values = [float(it.get("value", 1.0)) for it in items]
        total = sum(values) or 1.0

        chart_data = CategoryChartData()
        chart_data.categories = [str(i + 1) for i in range(n)]
        chart_data.add_series("Slices", values)
        gframe = slide.shapes.add_chart(
            XL_CHART_TYPE.DOUGHNUT,
            Inches(cx), Inches(cy), Inches(ring_d), Inches(ring_d),
            chart_data,
        )
        chart = gframe.chart
        chart.has_title = False
        chart.has_legend = False
        tints = _indigo_tints(n)
        try:
            series = chart.plots[0].series[0]
            for i in range(n):
                p = series.points[i]
                p.format.fill.solid()
                p.format.fill.fore_color.rgb = tints[i]
                p.format.line.color.rgb = TRENCH
                p.format.line.width = Pt(1.5)
        except Exception:
            pass

        # Number labels drawn on top of the ring at each wedge's midpoint. Excel/PPT
        # donuts start at 12 o'clock (angle=0) and sweep clockwise. The label radius
        # is halfway between the inner and outer edges so it sits on the ring.
        center_x = cx + ring_d / 2
        center_y = cy + ring_d / 2
        outer_r = ring_d / 2
        inner_r = outer_r * 0.55       # python-pptx default hole is ~50%; 55% is the label ring
        label_r = (outer_r + inner_r) / 2
        cumulative = 0.0
        for i, v in enumerate(values):
            mid_frac = (cumulative + v / 2) / total
            angle = mid_frac * 2 * math.pi - math.pi / 2  # start at 12 o'clock
            lx = center_x + label_r * math.cos(angle)
            ly = center_y + label_r * math.sin(angle)
            cumulative += v
            box_w, box_h = 0.60, 0.45
            tb = slide.shapes.add_textbox(
                Inches(lx - box_w / 2), Inches(ly - box_h / 2),
                Inches(box_w), Inches(box_h),
            )
            tf = tb.text_frame
            tf.margin_left = tf.margin_right = Inches(0.02)
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            r = p.add_run()
            r.text = str(i + 1)
            r.font.name = FONT_BODY
            r.font.bold = True
            r.font.size = Pt(18)
            r.font.color.rgb = MIST if _is_dark(tints[i]) else TRENCH

    # ---- Office locations (brand "Our offices" slide) --------------------

    def add_office_locations(
        self,
        offices: Iterable[Dict[str, Any]],
        eyebrow: Optional[str] = None,
        map_image: Optional[str] = None,
    ) -> Slide:
        """
        Brand "Our offices" slide — eyebrow top-left, two columns of office entries
        on the left half, and a stylised Europe map with indigo-highlighted
        countries on the right.

        offices = [{"city": "Helsinki", "address": "Kalevankatu 30\\n00100 Helsinki"}, ...]
                  or [{"city": "Tampere", "note": "Coworking space"}, ...]

        The bundled map (`graphics/map-europe-finland-highlight.png`) is used by
        default. Pass `map_image` to override with a different map — useful when
        highlighting a different country set.
        """
        offices = list(offices)
        if not 1 <= len(offices) <= 8:
            raise ValueError("add_office_locations supports 1–8 offices")
        slide = self._add_slide("TITLE_ONLY")
        self._fill_placeholder(slide, 0, "")
        self._fill_placeholder(slide, 2, eyebrow or "Our offices")

        # Right-half map
        try:
            if map_image:
                map_path = Path(map_image)
                if not map_path.exists():
                    map_path = self._find_asset("graphics", map_image)
            else:
                map_path = self._find_asset("graphics", "map-europe-finland-highlight")
            # Position the map in the right 45% of the slide.
            slide.shapes.add_picture(
                str(map_path),
                Inches(5.40), Inches(0.20),
                height=Inches(5.20),
            )
        except FileNotFoundError:
            pass  # skip map if not available

        # Left-half office entries — 2 columns × up to 4 rows.
        col_lefts = [0.55, 3.00]
        row_top_0 = 1.10
        row_gap = 1.40
        for i, off in enumerate(offices[:8]):
            col = i % 2
            row = i // 2
            x = col_lefts[col]
            y = row_top_0 + row * row_gap
            # City (TRENCH bold, large)
            self.add_body_textbox(
                slide, off.get("city", ""),
                left=x, top=y, width=2.30, height=0.50,
                color=TRENCH, size_pt=22, bold=True,
            )
            # Address / note
            detail = off.get("address") or off.get("note") or ""
            if detail:
                self.add_body_textbox(
                    slide, detail,
                    left=x, top=y + 0.55, width=2.30, height=0.70,
                    color=TRENCH, size_pt=13,
                )
        return slide

    # ---- Principles split (brand "Safe Space Principles" slide) ---------

    def add_principles_split(
        self,
        title: str,
        body: str,
        statements: Iterable[str],
        eyebrow: Optional[str] = None,
    ) -> Slide:
        """
        Brand "Safe Space Principles"-style split. A narrow Mist column on the left
        carries the title + body paragraphs; a wide Trench panel on the right carries
        2–4 large italic MIST statements stacked top-to-bottom, plus a calm aquatic
        form cluster in the bottom-right corner.

        body        : multi-paragraph text (split on blank lines). Use **...** for bold
                      emphasis inline.
        statements  : ["Respect. Care. Think.", "Communicate. Don't assume.", ...]
        """
        statements = list(statements)
        if not 1 <= len(statements) <= 4:
            raise ValueError("add_principles_split supports 1–4 statements")
        slide = self._add_slide("TITLE_ONLY")
        self._fill_placeholder(slide, 0, "")
        # Eyebrow (optional)
        if eyebrow:
            self._fill_placeholder(slide, 2, eyebrow)
        else:
            self._fill_placeholder(slide, 2, "")

        # --- Right panel: full-height Trench block -------------------------
        panel_left, panel_top = 3.90, 0.0
        panel_w, panel_h = 10.0 - panel_left, 5.625
        panel = self.add_rect(slide, panel_left, panel_top, panel_w, panel_h, fill=TRENCH)
        self._send_to_back(panel)

        # Subtle indigo radial glow in the bottom-right corner.
        # The bundled "cluster-full-02" orb has a hard black square background that
        # clashes with the TRENCH panel (TRENCH is #1E0050, not #000000 — the black
        # corners stand out). Instead we draw a soft INDIGO-tint oval directly:
        # it stays visually calm, doesn't overlap the statements, and reads as an
        # accent glow rather than a competing illustration.
        glow_d = 2.10
        glow = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Inches(panel_left + panel_w - glow_d * 0.55),
            Inches(panel_h - glow_d * 0.60),
            Inches(glow_d), Inches(glow_d),
        )
        glow.fill.solid()
        glow.fill.fore_color.rgb = INDIGO
        glow.line.fill.background()
        self._apply_indigo_gradient(glow)
        self._send_to_back(glow)
        # Re-pin the Trench panel behind the glow so the glow sits ON the panel.
        self._send_to_back(panel)

        # Statements stacked vertically, centred on the right panel.
        n = len(statements)
        # Vertical budget: leave 0.55" padding top/bottom.
        top_pad, bot_pad = 0.55, 0.55
        avail = panel_h - top_pad - bot_pad
        gap = 0.35
        block_h = (avail - (n - 1) * gap) / n
        for i, stmt in enumerate(statements):
            y = top_pad + i * (block_h + gap)
            self.add_heading_textbox(
                slide, stmt,
                left=panel_left + 0.45, top=y,
                width=panel_w - 0.65, height=block_h,
                color=MIST, size_pt=36, italic=True,
            )

        # --- Left column: Mist background (already the page default) ------
        # Title (bold, large)
        self.add_body_textbox(
            slide, title,
            left=0.55, top=0.60, width=3.05, height=1.00,
            color=TRENCH, size_pt=26, bold=True,
        )
        # Body — multi-paragraph with **bold** markers.
        self._add_rich_paragraphs(
            slide, body,
            left=0.55, top=1.75, width=3.05, height=3.50,
            color=TRENCH, size_pt=12,
        )
        return slide

    # ---- Reference logos (brand "Our references" grid) -----------------

    def add_reference_logos(
        self,
        logos: Iterable[str | Path],
        title: Optional[str] = None,
        eyebrow: Optional[str] = None,
        cols: int = 4,
    ) -> Slide:
        """
        Grid of client/reference logos on a Mist background. Each entry is either an
        absolute path to a PNG/JPG, or an asset name resolvable under
        `assets/logos/reference/` or `assets/misc/` via the standard lookup.

        Up to 16 logos render comfortably on one slide. Logos are centred inside
        equally sized cells and scaled to fit — no stretching. If an asset can't be
        resolved it's skipped silently (no broken-image placeholders).

        The brand reference slide has no title and no eyebrow; both are optional.
        """
        logos = list(logos)
        if not logos:
            raise ValueError("add_reference_logos needs at least one logo")
        if cols < 1:
            cols = 4
        slide = self._add_slide("TITLE_ONLY")
        self._fill_placeholder(slide, 0, title or "")
        self._fill_placeholder(slide, 2, eyebrow or "")

        # Resolve logo paths — accept absolute paths, asset names, and Path objects.
        resolved: List[Path] = []
        for entry in logos:
            p = Path(entry) if not isinstance(entry, Path) else entry
            if p.is_absolute() and p.exists():
                resolved.append(p)
                continue
            # Try assets/logos/reference, then assets/misc, then assets/logos
            for subdir, prefix in [
                ("logos/reference", ""),
                ("logos/reference", "client-logo-"),
                ("misc", ""),
                ("misc", "client-logo-"),
                ("logos", "client-logo-"),
            ]:
                try:
                    resolved.append(self._find_asset(subdir, str(entry), prefix=prefix))
                    break
                except FileNotFoundError:
                    continue
        if not resolved:
            # Nothing to show — still return the slide so callers don't crash.
            return slide

        # Grid geometry — centred on the slide, generous cell padding.
        n = len(resolved)
        rows = math.ceil(n / cols)
        # Leave 0.70" side margin, 1.0" top/bottom depending on whether title/eyebrow live there.
        side_pad = 0.70
        top_pad = 1.05 if title else 0.70
        bot_pad = 0.70
        avail_w = 10.0 - 2 * side_pad
        avail_h = 5.625 - top_pad - bot_pad
        cell_w = avail_w / cols
        cell_h = avail_h / rows
        # Allow the logo to occupy ~70% of its cell to keep breathing room.
        logo_box_w = cell_w * 0.70
        logo_box_h = cell_h * 0.70

        for i, path in enumerate(resolved):
            row = i // cols
            col = i % cols
            # Cell centre
            cx = side_pad + (col + 0.5) * cell_w
            cy = top_pad + (row + 0.5) * cell_h
            # Use height-constrained add_picture so aspect ratio is preserved; then
            # re-centre if it's wider/narrower than the target box.
            try:
                pic = slide.shapes.add_picture(
                    str(path),
                    Inches(cx - logo_box_w / 2), Inches(cy - logo_box_h / 2),
                    height=Inches(logo_box_h),
                )
            except Exception:
                continue
            # If it ended up wider than the box, swap to width-constraint.
            if pic.width > Inches(logo_box_w):
                pic_ar = pic.width / pic.height
                new_w = Inches(logo_box_w)
                new_h = int(new_w / pic_ar)
                pic.width = new_w
                pic.height = new_h
            # Recentre
            pic.left = Inches(cx) - pic.width // 2
            pic.top = Inches(cy) - pic.height // 2
        return slide

    def add_thank_you(self, subtitle: str = URL) -> Slide:
        """Closing 'Thank you!' slide — dark indigo with aquatic forms."""
        slide = self._add_slide("TITLE_1")
        self._fill_placeholders(slide, {0: "Thank you!", 1: subtitle})
        return slide

    # ---- Rich paragraph helper (for multi-paragraph body with optional **bold**) ------

    def _add_rich_paragraphs(
        self,
        slide: Slide,
        text: str,
        left: float, top: float, width: float, height: float,
        color: RGBColor = TRENCH,
        size_pt: int = 14,
    ):
        """
        Add a textbox containing the `text`, split on `\\n\\n` into paragraphs. Within each
        paragraph, markers like `**word**` render bold. This is the standard Mearra body
        treatment — emphasis stays in the same colour, just bold.
        """
        tb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
        tf = tb.text_frame
        tf.word_wrap = True
        paragraphs = text.split("\n\n") if text else [""]
        for i, para_text in enumerate(paragraphs):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            # Split on **...** while keeping delimiters so we can bold them.
            tokens = _split_bold(para_text)
            for tok_text, is_bold in tokens:
                r = p.add_run()
                r.text = tok_text
                r.font.name = FONT_BODY
                r.font.size = Pt(size_pt)
                r.font.color.rgb = color
                r.font.bold = is_bold
        return tb

    # ---- Custom canvas slide ------------------------------------------------

    def add_custom_slide(
        self,
        title: Optional[str] = None,
        eyebrow: Optional[str] = None,
        layout: str = "TITLE_ONLY",
    ) -> Slide:
        """
        Return a slide that has correct brand chrome (logo, background, tagline if cover layout)
        but an otherwise-blank canvas, so the caller can add custom shapes/images/text.

        Use this when none of the named layouts fit but you still want the master's
        colors, fonts, logo, and tagline.
        """
        slide = self._add_slide(layout)
        if title is not None:
            self._fill_placeholder(slide, 0, title)
        if eyebrow is not None:
            self._fill_placeholder(slide, 2, eyebrow)
        return slide

    # ---- Asset helpers ------------------------------------------------------

    def add_image(
        self,
        slide: Slide,
        image_path: str | Path,
        left: float, top: float,
        width: Optional[float] = None, height: Optional[float] = None,
        fit: str = "stretch",
    ):
        """
        Add an image to a slide. Positions/sizes in inches.

        `fit` controls aspect behaviour when BOTH `width` and `height` are given:
        - "stretch" (default) : force the image to the exact width×height box (may distort).
        - "cover"             : fill the box by scaling + centre-cropping, preserves aspect.
                                Picture ends up exactly width×height on the slide, with
                                overflow trimmed equally from the edges of the original image.
        """
        if width is not None and height is not None and fit == "cover":
            # Compute crop fractions so the cropped image has the same aspect as the target,
            # then render at exactly the requested box. python-pptx crops are fractions of the
            # *original* image, applied before scaling, so the picture stays width×height.
            try:
                from PIL import Image
                with Image.open(str(image_path)) as im:
                    iw, ih = im.size
            except Exception:
                iw, ih = (1, 1)
            target_aspect = width / height
            image_aspect = iw / ih if ih else 1.0
            pic = slide.shapes.add_picture(
                str(image_path),
                Inches(left), Inches(top),
                width=Inches(width), height=Inches(height),
            )
            if image_aspect > target_aspect:
                # Image is wider than target — crop equally left+right
                crop_total = 1 - target_aspect / image_aspect
                pic.crop_left = crop_total / 2
                pic.crop_right = crop_total / 2
            elif image_aspect < target_aspect:
                # Image is taller than target — crop equally top+bottom
                crop_total = 1 - image_aspect / target_aspect
                pic.crop_top = crop_total / 2
                pic.crop_bottom = crop_total / 2
            return pic
        # Default path: honour whatever (width, height) were passed verbatim.
        left_e, top_e = Inches(left), Inches(top)
        w_e = Inches(width) if width is not None else None
        h_e = Inches(height) if height is not None else None
        return slide.shapes.add_picture(str(image_path), left_e, top_e, width=w_e, height=h_e)

    def add_aquatic_form(
        self,
        slide: Slide,
        name: str,
        left: float = 6.5, top: float = 0.0,
        width: Optional[float] = 3.8, height: Optional[float] = None,
    ):
        """
        Drop an Aquatic Form illustration onto the slide.
        `name` matches files in assets/aquatic-forms/ (with or without extension).
        Default placement is the signature Mearra crop: right edge, extending beyond the frame.
        """
        path = self._find_asset("aquatic-forms", name)
        return self.add_image(slide, path, left, top, width=width, height=height)

    def add_icon(
        self,
        slide: Slide,
        name: str,
        left: float, top: float,
        size: float = 0.7,
    ):
        """Drop a named icon from assets/icons/. Icons are square; `size` is width in inches."""
        path = self._find_asset("icons", name, prefix="icon-")
        return self.add_image(slide, path, left, top, width=size, height=size)

    def add_portrait(
        self,
        slide: Slide,
        name: str,
        left: float = 6.5, top: float = 0.8,
        width: float = 3.0, height: Optional[float] = None,
    ):
        """Place a portrait from assets/portraits/ (person photos cut out for collage)."""
        path = self._find_asset("portraits", name, prefix="person-")
        return self.add_image(slide, path, left, top, width=width, height=height)

    def add_mockup(
        self,
        slide: Slide,
        name: str,
        left: float, top: float,
        width: float = 5.0, height: Optional[float] = None,
    ):
        """Place a device mockup (laptop, phone) from assets/mockups/ for product demo slides."""
        path = self._find_asset("mockups", name)
        return self.add_image(slide, path, left, top, width=width, height=height)

    def add_graphic(
        self,
        slide: Slide,
        name: str,
        left: float, top: float,
        width: float = 4.0, height: Optional[float] = None,
    ):
        """Place a reusable graphic (map, chart decoration) from assets/graphics/."""
        path = self._find_asset("graphics", name)
        return self.add_image(slide, path, left, top, width=width, height=height)

    def add_logo(
        self,
        slide: Slide,
        variant: str = "symbol-dark",
        left: float = 8.8, top: float = 4.9,
        width: float = 0.9,
    ):
        """
        Place a Mearra logo. variant: 'symbol-dark', 'symbol-light', or 'wordmark-light'.
        Default placement is bottom-right content-slide position.
        Note: most layouts already place the logo via the master — use this only for
        custom slides that don't have the master chrome.
        """
        path = ASSETS / "logos" / f"mearra-{variant}.png"
        return self.add_image(slide, path, left, top, width=width)

    # ---- Brand-safe text helpers -------------------------------------------

    def add_heading_textbox(
        self,
        slide: Slide,
        text: str,
        left: float, top: float, width: float, height: float,
        color: RGBColor = TRENCH,
        size_pt: int = 40,
        italic: bool = True,
    ):
        """Add an on-brand heading textbox (Sofia Sans Extra Condensed, Black Italic by default)."""
        tb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        r = p.add_run()
        r.text = text
        r.font.name = FONT_HEADING
        r.font.bold = True
        r.font.italic = italic
        r.font.size = Pt(size_pt)
        r.font.color.rgb = color
        return tb

    def add_body_textbox(
        self,
        slide: Slide,
        text: str,
        left: float, top: float, width: float, height: float,
        color: RGBColor = TRENCH,
        size_pt: int = 14,
        bold: bool = False,
    ):
        """Add an on-brand body textbox (Sofia Sans Regular/Medium)."""
        tb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        r = p.add_run()
        r.text = text
        r.font.name = FONT_BODY
        r.font.bold = bold
        r.font.size = Pt(size_pt)
        r.font.color.rgb = color
        return tb

    def add_rect(
        self,
        slide: Slide,
        left: float, top: float, width: float, height: float,
        fill: RGBColor = INDIGO,
        no_line: bool = True,
    ):
        """Add a simple rectangle (useful for brand-color blocks and timeline bars)."""
        shp = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(left), Inches(top), Inches(width), Inches(height),
        )
        shp.fill.solid()
        shp.fill.fore_color.rgb = fill
        if no_line:
            shp.line.fill.background()
        return shp

    # ---- Asset lookup -------------------------------------------------------

    def _find_asset(self, subdir: str, name: str, prefix: str = "") -> Path:
        """
        Resolve an asset name to an absolute path.
        Matches name exactly, then tries with prefix, then fuzzy stem match.
        """
        folder = ASSETS / subdir
        # Exact
        for candidate in [name, f"{prefix}{name}", f"{name}.png", f"{prefix}{name}.png", f"{prefix}{name}.jpg"]:
            p = folder / candidate
            if p.exists():
                return p
        # Fuzzy: stem contains name
        for p in folder.iterdir():
            if name.lower() in p.stem.lower():
                return p
        raise FileNotFoundError(f"No asset matching {name!r} in {folder}")

    # ---- Save ---------------------------------------------------------------

    def save(self, out_path: str | Path) -> Path:
        out_path = Path(out_path)
        out_path.parent.mkdir(parents=True, exist_ok=True)
        self.prs.save(str(out_path))
        return out_path

    # ---- Introspection ------------------------------------------------------

    def list_layouts(self) -> List[str]:
        return sorted(self.layouts.keys())
